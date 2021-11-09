import sys
from pprint import pprint
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from pptx.dml.color import _NoneColor

# CONFIG Variables
OUR_BACKGROUND_TEXTBOX_RGB = [ 68 , 114 , 196 ]
OUR_BACKGROUND_TEXTBOX_HEX = "0070C0"
MAXIMUM_TEXTBOX_HEIGHT = 503732
EXPORT_WIDTH = 1920
EXPORT_HEIGHT = 1080

# GLOBAL Variables
PARSED_TEXT_BOXES = []
PARSED_BLANK_BOXES = []

# https://stackoverflow.com/a/24290026
def enumerate2( xs , start=0 , step=1):
	for x in xs:
		yield ( start , x )
		start += step

def is_rectangular( shape ):
	if shape.width > shape.height:
		return True
	return False

def is_accepted_rectangular_form( shape ):
	# if is_rectangular( shape ) == False:
	# 	return False
	# if shape.height > MAXIMUM_TEXTBOX_HEIGHT:
	# 	return False
	return True

def get_fill_type( fill ):
	if hasattr( fill.type , "real" ) == False:
		return False
	fill_type = fill.type
	if fill_type == 1:
		return "solid"
	elif fill_type == 2:
		return "patterned"
	elif fill_type == 3:
		return "gradient"
	elif fill_type == 4:
		return "textured"
	elif fill_type == 5:
		return "background"
	elif fill_type == 6:
		return "picture"
	return False

def get_fill_color( shape ):
	if hasattr( shape , "fill" ) == False:
		return False
	fill_type = get_fill_type( shape.fill )
	if fill_type == False:
		return False
	if fill_type != "solid":
		return False
	color = shape.fill.fore_color
	result = {}
	if not isinstance( color._color , _NoneColor ):
		if hasattr( color, "rgb" ):
			result[ "rgb" ] = color.rgb
		if hasattr( color , "brightness" ):
			result[ "brightness" ] = color.brightness
		if hasattr( color , "color_type" ):
			result[ "color_type" ] = color.color_type
		if hasattr( color , "theme_color" ):
			result[ "theme_color" ] = color.theme_color
		if hasattr( color , "theme_color" ):
			result[ "theme_color" ] = "{}".format( color.theme_color )
	return result

def validate_correct_color( rgb_object ):
	if str( rgb_object ) == OUR_BACKGROUND_TEXTBOX_HEX:
		return True
	elif isinstance( rgb_object , tuple ):
		if rgb_object[ 0 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 0 ]:
			if rgb_object[ 1 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 1 ]:
				if rgb_object[ 2 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 2 ]:
					return True
	return False

def validate_text_box( shape ):
	# 1.) Validate Background Color Matches "our" Textboxes
	# NOTE: YOU HAVE TO USE SHAPE FILL BUTTON , idk man
	fill_color = get_fill_color( shape )
	if fill_color == False:
		return False
	if "rgb" not in fill_color:
		return False
	if validate_correct_color( fill_color[ "rgb" ] ) == False:
		return False

	# 2.) Validate There is Text In the Box
	if hasattr( shape , "text_frame" ) == False:
		return False
	if len( shape.text_frame.text ) < 1:
		return False

	# 3.) Validate Shape is "rEcTanGULar" and appropriate size
	# doesn't do a very good job, but 99% it works
	if is_accepted_rectangular_form( shape ) == False:
		return False
	return True

def validate_blank_box( shape ):
	# 1.) Validate Background Color Matches "our" Textboxes
	fill_color = get_fill_color( shape )
	if fill_color == False:
		return False
	if "rgb" not in fill_color:
		return False
	if validate_correct_color( fill_color[ "rgb" ] ) == False:
		return False

	# 2.) Validate There is Text In the Box
	if hasattr( shape , "text_frame" ) == False:
		return False
	if len( shape.text_frame.text ) > 0:
		return False

	# 3.) Validate Shape is "rEcTanGULar" and appropriate size
	# doesn't do a very good job, but 99% it works
	if is_accepted_rectangular_form( shape ) == False:
		return False
	return True

# This assumes that every textbox has a unique positon , across all slides
def find_matching_blank_box( text_box ):
	for index , blank_box in enumerate( PARSED_BLANK_BOXES ):
		if text_box[ "top" ] == blank_box[ "top" ]:
			if text_box[ "left" ] == blank_box[ "left" ]:
				if text_box[ "height" ] == blank_box[ "height" ]:
					if text_box[ "width" ] == blank_box[ "width" ]:
						return blank_box
	return False

# https://python-pptx.readthedocs.io/en/latest/api/shapes.html
# https://stackoverflow.com/questions/54692768/python-pptx-read-font-color
# https://github.com/scanny/python-pptx/blob/master/pptx/enum/dml.py
# https://github.com/BlenderCN/Learnbgame/blob/b6aabd1f331dc08915678f6e8a7706f8ac02ef46/All_In_One/addons/Presentation-master/powerPointConverter.py
# https://github.com/scanny/python-pptx/issues/534

# ok , so we don't have to do this bull shit validation of matching text_filled boxes with blank boxes
# we supposedly have a powerpoint library , just duplicate the slides without text while we are here
# easier for everyone

# https://python-pptx.readthedocs.io/en/latest/api/presentation.html#presentation-function
p = Presentation( sys.argv[ 1 ] )
print( p.slide_width )
print( p.slide_height )

# print( dir( p.slides[0].slide_layout ) )
# print( dir( p.slides[0].slide_layout.element ) )
# print( dir( p.slides[0].slide_layout.slide_master ) )
for slide_index , slide in enumerate( p.slides ):
	for shape_index , shape in enumerate( slide.shapes ):
		if validate_text_box( shape ):
			PARSED_TEXT_BOXES.append({
				"text": shape.text_frame.text ,
				"slide_number": ( slide_index + 1 ) ,
				"element_number": ( shape_index + 1 ) ,
				"top": shape.top ,
				"left": shape.left ,
				"height": shape.height ,
				"width": shape.width ,
			})
		elif validate_blank_box( shape ):
			PARSED_BLANK_BOXES.append({
				"slide_number": ( slide_index + 1 ) ,
				"element_number": ( shape_index + 1 ) ,
				"top": shape.top ,
				"left": shape.left ,
				"height": shape.height ,
				"width": shape.width ,
			})
if len( PARSED_TEXT_BOXES ) != len( PARSED_BLANK_BOXES ):
	print( "Was not able to parse the same number of filled and blank text boxes" )
	print( "you probably didn't use the 'shape fill' button for all the colors" )
	sys.exit( 1 )

matching_boxes = []
for index , text_box in enumerate( PARSED_TEXT_BOXES ):
	matching_blank_box = find_matching_blank_box( text_box )
	if matching_blank_box == False:
		print( "shapes didn't line up perfectly , you moved one on accident" )
		sys.exit( 1 )
	matching_boxes.append( [ text_box , matching_blank_box ] )

pprint( matching_boxes )




## Now is the hard part of tyring to figure out the coordinate space of the exported image , and then matching it to whatever the fuck an image-map is in HTML
## If you can do this , then we can automate the whole image mapping stuff

# https://39363.org/IMAGE_BUCKET/1636437240783-90278400.jpeg
# <map name="image-map">
# 	<area target="" alt="Cervical Enlargement" title="Cervical Enlargement" href="" coords="469,251,817,200" shape="rect">
# 	<area target="" alt="Lumbosacral Englargement" title="Lumbosacral Englargement" href="" coords="355,661,774,614" shape="rect">
# 	<area target="" alt="Conus Medullaris" title="Conus Medullaris" href="" coords="1144,673,1439,621" shape="rect">
# 	<area target="" alt="Cauda Equina" title="Cauda Equina" href="" coords="1144,780,1388,732" shape="rect">
# 	<area target="" alt="Filum Terminale" title="Filum Terminale" href="" coords="1144,1007,1439,956" shape="rect">
# </map>

# HTML Map
# function translate_raw_rectangle_coordinates_to_ordered( raw_input_coordinates ) {
# 	let input = raw_input_coordinates.split( "," ).map( x => parseInt( x ) );
# 	input[ 0 ] = Math.round( ( input[ 0 ] * width_scale_percentage ) );
# 	input[ 1 ] = Math.round( ( input[ 1 ] * height_scale_percentage ) );
# 	input[ 2 ] = Math.round( ( input[ 2 ] * width_scale_percentage ) );
# 	input[ 3 ] = Math.round( ( input[ 3 ] * height_scale_percentage ) );
# 	let ordered;
# 	if ( input[ 0 ] > input[ 2 ] ) {
# 		ordered = [ input[ 2 ] , input[ 3 ] , input[ 0 ] , input[ 1 ] ];
# 	}
# 	else {
# 		ordered = input;
# 	}
# 	let x_size = ( ordered[ 2 ] - ordered[ 0 ] );
# 	let y_size = ( ordered[ 3 ] - ordered[ 1 ] );
# 	// let y_size = Math.abs( ordered[ 3 ] - ordered[ 1 ] );
# 	let final = [ ordered[ 0 ] , ordered[ 1 ] , x_size , y_size ];
# 	return final;
# }

# Slide Coordinates
# Left = -9939
# Top = -13403
# Bottom = 6569765 +- 298174
# Right = 10336693 +- 1855307

example = {
	"filum_terminale": {
		"html_map_raw": "1144,1007,1439,956" ,
		"power_point": "7268814,6089297,1855307,298174"
	}
}

 # [{'element_number': 6,
 #   'height': 298174,
 #   'left': 7268814,
 #   'slide_number': 1,
 #   'text': 'Filum Terminale',
 #   'top': 6089297,
 #   'width': 1855307},


# [[{'element_number': 2,
#    'height': 298174,
#    'left': 2987749,
#    'slide_number': 1,
#    'text': 'Cervical Enlargement',
#    'top': 1278683,
#    'width': 2180599},
#   {'element_number': 2,
#    'height': 298174,
#    'left': 2987749,
#    'slide_number': 2,
#    'top': 1278683,
#    'width': 2180599}],
#  [{'element_number': 3,
#    'height': 298174,
#    'left': 2256184,
#    'slide_number': 1,
#    'text': 'Lumbosacral Enlargement',
#    'top': 3919254,
#    'width': 2647122},
#   {'element_number': 3,
#    'height': 298174,
#    'left': 2256184,
#    'slide_number': 2,
#    'top': 3919254,
#    'width': 2647122}],
#  [{'element_number': 4,
#    'height': 298174,
#    'left': 7268815,
#    'slide_number': 1,
#    'text': 'Cauda Equina',
#    'top': 4658063,
#    'width': 1517376},
#   {'element_number': 4,
#    'height': 298174,
#    'left': 7268815,
#    'slide_number': 2,
#    'top': 4658063,
#    'width': 1517376}],
#  [{'element_number': 5,
#    'height': 298174,
#    'left': 7268814,
#    'slide_number': 1,
#    'text': 'Conus Medullaris',
#    'top': 3959010,
#    'width': 1855307},
#   {'element_number': 5,
#    'height': 298174,
#    'left': 7268814,
#    'slide_number': 2,
#    'top': 3959010,
#    'width': 1855307}],
#  [{'element_number': 6,
#    'height': 298174,
#    'left': 7268814,
#    'slide_number': 1,
#    'text': 'Filum Terminale',
#    'top': 6089297,
#    'width': 1855307},
#   {'element_number': 6,
#    'height': 298174,
#    'left': 7268814,
#    'slide_number': 2,
#    'top': 6089297,
#    'width': 1855307}]]


















