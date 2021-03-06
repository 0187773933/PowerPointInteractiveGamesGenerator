import sys
from pprint import pprint
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches , Pt
from pptx.oxml import parse_xml
from pptx.dml.color import _NoneColor , RGBColor

# https://github.com/scanny/python-pptx/blob/62f7f6211c29150245fd8b383c5bde69534f48dd/pptx/util.py#L16
_EMUS_PER_INCH = 914400
_EMUS_PER_CENTIPOINT = 127
_EMUS_PER_CM = 360000 # centimeters
_EMUS_PER_MM = 36000 # millimeters
_EMUS_PER_PT = 12700

# https://github.com/scanny/python-pptx/blob/62f7f6211c29150245fd8b383c5bde69534f48dd/pptx/shapes/base.py#L93
# All in English Metric Units (EMU)
# Height = Integer distance between top and bottom extents of shape
# Width = Integer distance between left and right extents of shape
# Left = Integer distance of the left edge of this shape from the left edge of the slide
# Top = Integer distance of the top edge of this shape from the top edge of the slide

# A convenient method for calculating scale is to divide a Length object by an equivalent count of local coordinate units
# e.g. scale = Inches(1)/1000 for 1000 local units per inch.

# CONFIG Variables
OUR_BACKGROUND_TEXTBOX_RGB = [ 68 , 114 , 196 ]
OUR_BACKGROUND_TEXTBOX_RGB_TUPLE = ( 68 , 114 , 196 )
OUR_BACKGROUND_TEXTBOX_HEX = "0070C0"
MAXIMUM_TEXTBOX_HEIGHT = 503732
SLIDE_MASTER_WIDTH = 12192000
SLIDE_MASTER_HEIGHT = 6858000
SLIDE_WIDTH_IN_INCHES = ( SLIDE_MASTER_WIDTH / _EMUS_PER_INCH )
SLIDE_HEIGHT_IN_INCHES = ( SLIDE_MASTER_HEIGHT / _EMUS_PER_INCH )
#print( SLIDE_WIDTH_IN_INCHES , SLIDE_HEIGHT_IN_INCHES )
SLIDE_MIDPOINT_IN_INCHES = [ ( SLIDE_WIDTH_IN_INCHES / 2 ) , ( SLIDE_HEIGHT_IN_INCHES / 2 ) ]
# https://github.com/scanny/python-pptx/blob/master/pptx/util.py#L69

EXPORT_IMAGE_DPI = 144
EXPORTED_WIDTH = 1920
EXPORTED_HEIGHT = 1080
EXPORTED_WIDTH_IN_INCHES = ( EXPORTED_WIDTH / EXPORT_IMAGE_DPI )
EXPORTED_HEIGTH_IN_INCHES = ( EXPORTED_HEIGHT / EXPORT_IMAGE_DPI )
#print( EXPORTED_WIDTH_IN_INCHES , EXPORTED_HEIGTH_IN_INCHES )
# identify -format "%s,%y\n" TestGen/Slide1.jpeg
# exiftool -p '$XResolution,$YResolution' TestGen/Slide1.jpeg

# GLOBAL Variables
PARSED_TEXT_BOXES = []
PARSED_BLANK_BOXES = []

# https://stackoverflow.com/a/24290026
def enumerate2( xs , start=0 , step=1):
	for x in xs:
		yield ( start , x )
		start += step

def is_rectangular( shape ):
	if shape.width > shape.height:
		return True
	return False

def is_accepted_rectangular_form( shape ):
	# if is_rectangular( shape ) == False:
	# 	return False
	# if shape.height > MAXIMUM_TEXTBOX_HEIGHT:
	# 	return False
	return True

def get_fill_type( fill ):
	if hasattr( fill.type , "real" ) == False:
		return False
	fill_type = fill.type
	if fill_type == 1:
		return "solid"
	elif fill_type == 2:
		return "patterned"
	elif fill_type == 3:
		return "gradient"
	elif fill_type == 4:
		return "textured"
	elif fill_type == 5:
		return "background"
	elif fill_type == 6:
		return "picture"
	return False

def get_fill_color( shape ):
	if hasattr( shape , "fill" ) == False:
		return False
	fill_type = get_fill_type( shape.fill )
	if fill_type == False:
		return False
	if fill_type != "solid":
		return False
	color = shape.fill.fore_color
	result = {}
	if not isinstance( color._color , _NoneColor ):
		if hasattr( color, "rgb" ):
			result[ "rgb" ] = color.rgb
		if hasattr( color , "brightness" ):
			result[ "brightness" ] = color.brightness
		if hasattr( color , "color_type" ):
			result[ "color_type" ] = color.color_type
		if hasattr( color , "theme_color" ):
			result[ "theme_color" ] = color.theme_color
		if hasattr( color , "theme_color" ):
			result[ "theme_color" ] = "{}".format( color.theme_color )
	return result

def validate_correct_color( rgb_object ):
	if str( rgb_object ) == OUR_BACKGROUND_TEXTBOX_HEX:
		return True
	elif isinstance( rgb_object , tuple ):
		if rgb_object[ 0 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 0 ]:
			if rgb_object[ 1 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 1 ]:
				if rgb_object[ 2 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 2 ]:
					return True
	return False

def validate_text_box( shape ):
	# 1.) Validate Background Color Matches "our" Textboxes
	# NOTE: YOU HAVE TO USE SHAPE FILL BUTTON , idk man
	fill_color = get_fill_color( shape )
	if fill_color == False:
		return False
	if "rgb" not in fill_color:
		return False
	if validate_correct_color( fill_color[ "rgb" ] ) == False:
		return False

	# 2.) Validate There is Text In the Box
	if hasattr( shape , "text_frame" ) == False:
		return False
	if len( shape.text_frame.text ) < 1:
		return False

	# 3.) Validate Shape is "rEcTanGULar" and appropriate size
	# doesn't do a very good job, but 99% it works
	if is_accepted_rectangular_form( shape ) == False:
		return False
	return True

def validate_blank_box( shape ):
	# 1.) Validate Background Color Matches "our" Textboxes
	fill_color = get_fill_color( shape )
	if fill_color == False:
		return False
	if "rgb" not in fill_color:
		return False
	if validate_correct_color( fill_color[ "rgb" ] ) == False:
		return False

	# 2.) Validate There is Text In the Box
	if hasattr( shape , "text_frame" ) == False:
		return False
	if len( shape.text_frame.text ) > 0:
		return False

	# 3.) Validate Shape is "rEcTanGULar" and appropriate size
	# doesn't do a very good job, but 99% it works
	if is_accepted_rectangular_form( shape ) == False:
		return False
	return True

# This assumes that every textbox has a unique positon , across all slides
def find_matching_blank_box( text_box ):
	for index , blank_box in enumerate( PARSED_BLANK_BOXES ):
		if text_box[ "top" ] == blank_box[ "top" ]:
			if text_box[ "left" ] == blank_box[ "left" ]:
				if text_box[ "height" ] == blank_box[ "height" ]:
					if text_box[ "width" ] == blank_box[ "width" ]:
						return blank_box
	return False

def get_shape_objects( input_path ):
	p = Presentation( input_path )
	print( p.slide_width )
	print( p.slide_height )

	# print( dir( p.slides[0].slide_layout ) )
	# print( dir( p.slides[0].slide_layout.element ) )
	# print( dir( p.slides[0].slide_layout.slide_master ) )
	for slide_index , slide in enumerate( p.slides ):
		for shape_index , shape in enumerate( slide.shapes ):
			if validate_text_box( shape ):
				PARSED_TEXT_BOXES.append({
					"text": shape.text_frame.text ,
					"slide_number": ( slide_index + 1 ) ,
					"element_number": ( shape_index + 1 ) ,
					"top": shape.top ,
					"left": shape.left ,
					"height": shape.height ,
					"width": shape.width ,
				})
			elif validate_blank_box( shape ):
				PARSED_BLANK_BOXES.append({
					"slide_number": ( slide_index + 1 ) ,
					"element_number": ( shape_index + 1 ) ,
					"top": shape.top ,
					"left": shape.left ,
					"height": shape.height ,
					"width": shape.width ,
				})
	if len( PARSED_TEXT_BOXES ) != len( PARSED_BLANK_BOXES ):
		print( "Was not able to parse the same number of filled and blank text boxes" )
		print( "you probably didn't use the 'shape fill' button for all the colors" )
		sys.exit( 1 )

	matching_boxes = []
	for index , text_box in enumerate( PARSED_TEXT_BOXES ):
		matching_blank_box = find_matching_blank_box( text_box )
		if matching_blank_box == False:
			print( "shapes didn't line up perfectly , you moved one on accident" )
			sys.exit( 1 )
		matching_boxes.append( [ text_box , matching_blank_box ] )

	pprint( matching_boxes )
	return matching_boxes

def convert_power_point_raw_to_inches( raw ):
	return [ ( x / _EMUS_PER_INCH ) for x in raw ]

# https://39363.org/IMAGE_BUCKET/1636480125222-87384442.jpeg
# Based on X-Values , the two coordinates must be in order from lowest to highest
# So just reorder them if they are not , html figures it out , but we need them this way
# <map name="image-map">
# 	<area target="" alt="top_left" title="top_left" href="" coords="3,148,146,0" shape="rect">
# 	<area target="" alt="top_right" title="top_right" href="" coords="1774,146,1918,2" shape="rect">
# 	<area target="" alt="bottom_left" title="bottom_left" href="" coords="-2,1078,146,934" shape="rect">
# 	<area target="" alt="bottom_right" title="bottom_right" href="" coords="1776,1077,1918,932" shape="rect">
# 	<area target="" alt="test" title="test" href="" coords="960,685,1106,537" shape="rect">
# </map>

## HTML Logic
# // Cross
# // draw_line( 0 , 0 , scaled_x , scaled_y ); // top left corner to bottom right corner
# // draw_line( 0 , scaled_y , scaled_x , 0 ); // bottom left corner to top right corner

# // Middle
# let midpoint_x = ( scaled_x / 2 );
# let midpoint_y = ( scaled_y / 2 );
# // draw_line( midpoint_x , 0 , midpoint_x , scaled_y ); // vertical across
# // draw_line( 0 , midpoint_y , scaled_x , midpoint_y ); // horizontal across

## PowerPoint Logic
# add_shape( left , top , width , height )
# All in English Metric Units (EMU)
# Left = Integer distance of the left edge of this shape from the left edge of the slide
# Top = Integer distance of the top edge of this shape from the top edge of the slide
# Width = Integer distance between left and right extents of shape
# Height = Integer distance between top and bottom extents of shape
def compute_image_map( input_object ):
	print( "" )
	inches = convert_power_point_raw_to_inches( input_object[ "power_point" ][ "raw" ] )
	translated_x1 = inches[ 0 ] * EXPORT_IMAGE_DPI
	translated_y1 = inches[ 1 ] * EXPORT_IMAGE_DPI
	translated_x2 = inches[ 2 ] * EXPORT_IMAGE_DPI
	translated_y2 = inches[ 3 ] * EXPORT_IMAGE_DPI
	input_object[ "power_point" ][ "inches" ] = inches
	x1 = translated_x1
	y1 = ( translated_y1 + translated_y2 )
	x2 = ( translated_x1 + translated_x2 )
	y2 = translated_y1
	input_object[ "image_map" ][ "translated" ] = [ x1 , y1 , x2 , y2 ]
	print( input_object[ "power_point" ][ "inches" ] )
	print( input_object[ "image_map" ][ "translated" ] )
	return input_object

shape_objects = get_shape_objects( sys.argv[ 1 ] )
for index , shape_object in enumerate( shape_objects ):
	print( f'{shape_object[ 0 ][ "text" ]} ===' )
	result = compute_image_map({
		"power_point": {
			"raw": [ shape_object[ 0 ][ "left" ] , shape_object[ 0 ][ "top" ] , shape_object[ 0 ][ "width" ] , shape_object[ 0 ][ "height" ] ]
		} ,
		"image_map": {}
	})






















