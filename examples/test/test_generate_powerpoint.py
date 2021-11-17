import sys
from pprint import pprint
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches , Pt
from pptx.oxml import parse_xml
from pptx.dml.color import _NoneColor , RGBColor

# https://github.com/scanny/python-pptx/blob/62f7f6211c29150245fd8b383c5bde69534f48dd/pptx/util.py#L16
_EMUS_PER_INCH = 914400
_EMUS_PER_CENTIPOINT = 127
_EMUS_PER_CM = 360000 # centimeters
_EMUS_PER_MM = 36000 # millimeters
_EMUS_PER_PT = 12700

# https://github.com/scanny/python-pptx/blob/62f7f6211c29150245fd8b383c5bde69534f48dd/pptx/shapes/base.py#L93
# All in English Metric Units (EMU)
# Height = Integer distance between top and bottom extents of shape
# Width = Integer distance between left and right extents of shape
# Left = Integer distance of the left edge of this shape from the left edge of the slide
# Top = Integer distance of the top edge of this shape from the top edge of the slide

# A convenient method for calculating scale is to divide a Length object by an equivalent count of local coordinate units
# e.g. scale = Inches(1)/1000 for 1000 local units per inch.

# CONFIG Variables
OUR_BACKGROUND_TEXTBOX_RGB = [ 68 , 114 , 196 ]
OUR_BACKGROUND_TEXTBOX_RGB_TUPLE = ( 68 , 114 , 196 )
OUR_BACKGROUND_TEXTBOX_HEX = "0070C0"
MAXIMUM_TEXTBOX_HEIGHT = 503732
SLIDE_MASTER_WIDTH = 12192000
SLIDE_MASTER_HEIGHT = 6858000
SLIDE_WIDTH_IN_INCHES = ( SLIDE_MASTER_WIDTH / _EMUS_PER_INCH )
SLIDE_HEIGHT_IN_INCHES = ( SLIDE_MASTER_HEIGHT / _EMUS_PER_INCH )
#print( SLIDE_WIDTH_IN_INCHES , SLIDE_HEIGHT_IN_INCHES )
SLIDE_MIDPOINT_IN_INCHES = [ ( SLIDE_WIDTH_IN_INCHES / 2 ) , ( SLIDE_HEIGHT_IN_INCHES / 2 ) ]
# https://github.com/scanny/python-pptx/blob/master/pptx/util.py#L69

EXPORT_IMAGE_DPI = 144
EXPORTED_WIDTH = 1920
EXPORTED_HEIGHT = 1080
EXPORTED_WIDTH_IN_INCHES = ( EXPORTED_WIDTH / EXPORT_IMAGE_DPI )
EXPORTED_HEIGTH_IN_INCHES = ( EXPORTED_HEIGHT / EXPORT_IMAGE_DPI )
#print( EXPORTED_WIDTH_IN_INCHES , EXPORTED_HEIGTH_IN_INCHES )
# identify -format "%s,%y\n" TestGen/Slide1.jpeg
# exiftool -p '$XResolution,$YResolution' TestGen/Slide1.jpeg


# Left , Top , Width , Height
# ## Add Shape Example With Embedded Text
def add_our_custom_textbox_to_slide_using_inches( slide , text , x , y , width , height ):
	location = [ Inches( x ) , Inches( y ) , Inches( width ) , Inches( height ) ]
	#print( location )
	added_shape = slide.shapes.add_shape( MSO_SHAPE.ROUNDED_RECTANGLE , *location )
	paragraph = added_shape.text_frame.add_paragraph()
	paragraph.text = text
	added_shape.fill.solid()
	added_shape.fill.fore_color.rgb = RGBColor.from_string( OUR_BACKGROUND_TEXTBOX_HEX )

def add_textbox_example( slide ):
	# add_textbox( left , top , width , height )
	textbox = first_slide.shapes.add_textbox( Inches( 3 ) , Inches( 1.5 ) , Inches( 3 ) , Inches( 1 ) )
	textframe = textbox.text_frame
	paragraph = textframe.add_paragraph()
	paragraph.text = "This is a paragraph in the slide!"

# https://www.tutorialspoint.com/how-to-create-powerpoint-files-using-python
p = Presentation()
p.slide_width = SLIDE_MASTER_WIDTH
p.slide_height = SLIDE_MASTER_HEIGHT
blank_layout = p.slide_layouts[ 6 ]
first_slide = p.slides.add_slide( blank_layout )


x_Textbox_width = 1.0
x_Textbox_height = 1.0
# Top Left Corner
add_our_custom_textbox_to_slide_using_inches( first_slide , "-1 , 1" , 0 , 0 , x_Textbox_width , x_Textbox_height )
# Top Right Corner
add_our_custom_textbox_to_slide_using_inches( first_slide , "1 , 1" , ( SLIDE_WIDTH_IN_INCHES - x_Textbox_width ) , 0 , x_Textbox_width , x_Textbox_height )
# Bottom Left Corner
add_our_custom_textbox_to_slide_using_inches( first_slide , "-1 , -1" , 0 , ( SLIDE_HEIGHT_IN_INCHES - x_Textbox_height ) , x_Textbox_width , x_Textbox_height )
# Bottom Right Corner
add_our_custom_textbox_to_slide_using_inches( first_slide , "-1 , 1" , ( SLIDE_WIDTH_IN_INCHES - x_Textbox_width ) , ( SLIDE_HEIGHT_IN_INCHES - x_Textbox_height ) , x_Textbox_width , x_Textbox_height )

# Center - Expanding - Down and Right
add_our_custom_textbox_to_slide_using_inches( first_slide , "test" , SLIDE_MIDPOINT_IN_INCHES[ 0 ] , SLIDE_MIDPOINT_IN_INCHES[ 1 ] , x_Textbox_width , x_Textbox_height )

# https://39363.org/IMAGE_BUCKET/1636480125222-87384442.jpeg
# Based on X-Values , the two coordinates must be in order from lowest to highest
# So just reorder them if they are not , html figures it out , but we need them this way
# <map name="image-map">
# 	<area target="" alt="top_left" title="top_left" href="" coords="3,148,146,0" shape="rect">
# 	<area target="" alt="top_right" title="top_right" href="" coords="1774,146,1918,2" shape="rect">
# 	<area target="" alt="bottom_left" title="bottom_left" href="" coords="-2,1078,146,934" shape="rect">
# 	<area target="" alt="bottom_right" title="bottom_right" href="" coords="1776,1077,1918,932" shape="rect">
# 	<area target="" alt="test" title="test" href="" coords="960,685,1106,537" shape="rect">
# </map>

example = {
	"top_left": {
		"power_point": {
			"raw": [ 0 , 0 , 914400 , 914400 ] ,
			"inches": [  ]
		} ,
		"image_map": {
			"raw": [ 0 , 148 , 146 , 0 ] ,
			"inches": [  ]
		}
	} ,
	"top_right": {
		"power_point": {
			"raw": [ 11277600 , 0 , 914400 , 914400 ] ,
			"inches": [ ]
		} ,
		"image_map": {
			"raw": [ 1774 , 146 , 1918 , 0 ] ,
			"inches": [  ]
		}
	} ,
	"bottom_left": {
		"power_point": {
			"raw": [ 0 , 5943600 , 914400 , 914400 ] ,
			"inches": [ ]
		} ,
		"image_map": {
			"raw": [ 0 , 1078 , 146 , 934 ] ,
			"inches": [  ]
		}
	} ,
	"bottom_right": {
		"power_point": {
			"raw": [ 11277600 , 5943600 , 914400 , 914400 ] ,
			"inches": [ ]
		} ,
		"image_map": {
			"raw": [ 1776 , 1077 , 1918 , 932 ] ,
			"inches": [  ]
		}
	} ,
	"test": {
		"power_point": {
			"raw": [ 6096000 , 3429000 , 914400 , 914400 ] ,
			"inches": [  ]
		} ,
		"image_map": {
			"raw": [ 960 , 685 , 1106 , 537 ] ,
			"inches": [  ]
		}
	}
}
## HTML Logic
# // Cross
# // draw_line( 0 , 0 , scaled_x , scaled_y ); // top left corner to bottom right corner
# // draw_line( 0 , scaled_y , scaled_x , 0 ); // bottom left corner to top right corner

# // Middle
# let midpoint_x = ( scaled_x / 2 );
# let midpoint_y = ( scaled_y / 2 );
# // draw_line( midpoint_x , 0 , midpoint_x , scaled_y ); // vertical across
# // draw_line( 0 , midpoint_y , scaled_x , midpoint_y ); // horizontal across

## PowerPoint Logic
# add_shape( left , top , width , height )
# All in English Metric Units (EMU)
# Left = Integer distance of the left edge of this shape from the left edge of the slide
# Top = Integer distance of the top edge of this shape from the top edge of the slide
# Width = Integer distance between left and right extents of shape
# Height = Integer distance between top and bottom extents of shape

def convert_power_point_raw_to_inches( raw ):
	return [ ( x / _EMUS_PER_INCH ) for x in raw ]

def compute_image_map_for_top_left( top_left ):
	print( "" )
	print( "Top Left" )
	inches = convert_power_point_raw_to_inches( top_left[ "power_point" ][ "raw" ] )
	translated_x1 = inches[ 0 ] * EXPORT_IMAGE_DPI
	translated_y1 = inches[ 1 ] * EXPORT_IMAGE_DPI
	translated_x2 = inches[ 2 ] * EXPORT_IMAGE_DPI
	translated_y2 = inches[ 3 ] * EXPORT_IMAGE_DPI
	top_left[ "power_point" ][ "inches" ] = inches
	x1 = EXPORTED_WIDTH - ( EXPORTED_WIDTH - translated_x1 )
	y1 = translated_y2
	x2 = EXPORTED_WIDTH - ( EXPORTED_WIDTH - translated_x2 )
	y2 = translated_y1
	top_left[ "image_map" ][ "translated" ] = [ x1 , y1 , x2 , y2 ]
	print( top_left[ "power_point" ][ "inches" ] )
	print( top_left[ "image_map" ][ "raw" ] )
	print( top_left[ "image_map" ][ "translated" ] )
	return top_left

example[ "top_left" ] = compute_image_map_for_top_left( example[ "top_left" ] )

def compute_image_map_for_top_right( top_right ):
	print( "" )
	print( "Top Right" )
	inches = convert_power_point_raw_to_inches( top_right[ "power_point" ][ "raw" ] )
	translated_x1 = inches[ 0 ] * EXPORT_IMAGE_DPI
	translated_y1 = inches[ 1 ] * EXPORT_IMAGE_DPI
	translated_x2 = inches[ 2 ] * EXPORT_IMAGE_DPI
	translated_y2 = inches[ 3 ] * EXPORT_IMAGE_DPI
	print( "Translated === " , translated_x1 , translated_y1 , translated_x2 , translated_y2 )
	top_right[ "power_point" ][ "inches" ] = inches
	x1 = translated_x1
	y1 = translated_y2
	x2 = ( translated_x1 + translated_x2 )
	y2 = translated_y1
	top_right[ "image_map" ][ "translated" ] = [ x1 , y1 , x2 , y2 ]
	# print( top_right[ "power_point" ][ "inches" ] )
	print( "Goal === " , top_right[ "image_map" ][ "raw" ] )
	print( "Final === " , top_right[ "image_map" ][ "translated" ] )
	return top_right

example[ "top_right" ] = compute_image_map_for_top_right( example[ "top_right" ] )


def compute_image_map_for_bottom_left( bottom_left ):
	print( "" )
	print( "Bottom Left" )
	inches = convert_power_point_raw_to_inches( bottom_left[ "power_point" ][ "raw" ] )
	translated_x1 = inches[ 0 ] * EXPORT_IMAGE_DPI
	translated_y1 = inches[ 1 ] * EXPORT_IMAGE_DPI
	translated_x2 = inches[ 2 ] * EXPORT_IMAGE_DPI
	translated_y2 = inches[ 3 ] * EXPORT_IMAGE_DPI
	print( "Translated === " , translated_x1 , translated_y1 , translated_x2 , translated_y2 )
	bottom_left[ "power_point" ][ "inches" ] = inches
	x1 = translated_x1
	y1 = translated_y2 + translated_y1
	x2 = ( translated_x1 + translated_x2 )
	y2 = translated_y1
	bottom_left[ "image_map" ][ "translated" ] = [ x1 , y1 , x2 , y2 ]
	# print( bottom_left[ "power_point" ][ "inches" ] )
	print( "Goal === " , bottom_left[ "image_map" ][ "raw" ] )
	print( "Final === " , bottom_left[ "image_map" ][ "translated" ] )
	return bottom_left

example[ "bottom_left" ] = compute_image_map_for_bottom_left( example[ "bottom_left" ] )


## =================================================================================================
## =================================================================================================
## =================================================================================================
def compute_image_map_for_bottom_right( bottom_right ):
	print( "" )
	print( "Bottom Right" )
	inches = convert_power_point_raw_to_inches( bottom_right[ "power_point" ][ "raw" ] )
	translated_x1 = inches[ 0 ] * EXPORT_IMAGE_DPI
	translated_y1 = inches[ 1 ] * EXPORT_IMAGE_DPI
	translated_x2 = inches[ 2 ] * EXPORT_IMAGE_DPI
	translated_y2 = inches[ 3 ] * EXPORT_IMAGE_DPI
	bottom_right[ "power_point" ][ "inches" ] = inches
	x1 = translated_x1
	y1 = ( translated_y1 + translated_y2 )
	x2 = ( translated_x1 + translated_x2 )
	y2 = translated_y1
	bottom_right[ "image_map" ][ "translated" ] = [ x1 , y1 , x2 , y2 ]
	print( bottom_right[ "power_point" ][ "inches" ] )
	print( bottom_right[ "image_map" ][ "raw" ] )
	print( bottom_right[ "image_map" ][ "translated" ] )
	return bottom_right

example[ "bottom_right" ] = compute_image_map_for_bottom_right( example[ "bottom_right" ] )

def compute_image_map_for_test( test ):
	print( "" )
	print( "Test" )
	inches = convert_power_point_raw_to_inches( test[ "power_point" ][ "raw" ] )
	translated_x1 = inches[ 0 ] * EXPORT_IMAGE_DPI
	translated_y1 = inches[ 1 ] * EXPORT_IMAGE_DPI
	translated_x2 = inches[ 2 ] * EXPORT_IMAGE_DPI
	translated_y2 = inches[ 3 ] * EXPORT_IMAGE_DPI
	test[ "power_point" ][ "inches" ] = inches
	x1 = translated_x1
	y1 = translated_y1
	x2 = ( translated_x1 + translated_x2 )
	y2 = ( translated_y1 + translated_y2 )
	test[ "image_map" ][ "translated" ] = [ x1 , y1 , x2 , y2 ]
	print( test[ "power_point" ][ "inches" ] )
	print( test[ "image_map" ][ "raw" ] )
	print( test[ "image_map" ][ "translated" ] )
	return test
example[ "test" ] = compute_image_map_for_test( example[ "test" ] )

p.save( "TestGen.pptx" )