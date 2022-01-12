#!/usr/bin/env python3
import sys
import uuid
import time
import base64
from pathlib import Path
from pprint import pprint
from datetime import timedelta
import tempfile
import zipfile
import shutil
import requests

from sanic import Sanic
from sanic.response import html as sanic_html
from sanic.response import json as sanic_json
from sanic.response import file as sanic_file
from sanic.response import file_stream as sanic_file_stream
from sanic.response import stream as sanic_stream
from sanic.request import Request

from sanic_jwt_extended import JWT , jwt_required
from sanic_jwt_extended.tokens import Token

import utils
import compute_image_maps
import image_uploader
import interactive_notes_generator

from pptx import Presentation
from copy import deepcopy

DEFAULT_CONFIG = utils.read_json( sys.argv[1] )

def upload_file_to_imgur( file_path ):
	try:
		response = requests.post(
			"https://api.imgur.com/3/upload.json" ,
			headers={ "Authorization": f"Client-ID {DEFAULT_CONFIG[ 'image_upload_server_imgur_version' ][ 'imgur_client_id' ]}" } ,
			data={
				"image": base64.b64encode( open( str( file_path ) , "rb" ).read() ) ,
				"type": "base64" ,
				"name": file_path.stem ,
			}
		)
		response.raise_for_status()
		data = response.json()
		# pprint( data )
		if "data" not in data:
			return False
		if "link" not in data["data"]:
			return False
		return data["data"]["link"]
	except Exception as e:
		print( e )
		return False

app = Sanic( __name__ )

# https://sanic-jwt-extended.seonghyeon.dev/config_options.html
# https://pyjwt.readthedocs.io/en/latest/algorithms.html
with JWT.initialize( app ) as manager:
	manager.config.secret_key = DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ]
	manager.config.public_key = DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_public_key" ]
	manager.config.private_key = DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_private_key" ]
	manager.config.algorithm = "HS512"
	manager.config.jwt_header_key = "token"
	manager.config.jwt_header_prefix = "Bearer"
	manager.config.jwt_cookie = "ppt_interactive_image_upload_server"
	manager.config.access_token_expires = timedelta( weeks=( 52 * 10 ) )

# https://sanic-jwt-extended.seonghyeon.dev/api_docs/jwt.html#class-sanic_jwtextendedjwt
@app.route( "/login" , methods=[ "POST" ] )
async def login( request: Request ):
	username = request.json.get( "username" , "user" )
	access_token = JWT.create_access_token( identity=username )
	return sanic_json( dict( token=access_token ) , status=200 )

# https://sanic-jwt-extended.seonghyeon.dev/usage/basic.html#use-token-object
# https://github.com/NovemberOscar/Sanic-JWT-Extended/blob/0eb9282a4c21f6d9a81c3d3c2a4818353b79b429/sanic_jwt_extended/decorators.py#L11
@app.route( "/protected" , methods=[ "GET" ] )
@jwt_required
async def protected( request: Request , token: Token ):
	print( "accessed protected route" )
	return sanic_json( dict( identity=token.identity , type=token.type , raw_data=token.raw_data , exp=str( token.exp ) ) )


def update_config_with_form_key( request , form_key , config , *config_keys ):
	item = request.form.get( form_key )
	if item != None and len( item ) > 0:
		utils.set_nested_dictionary_value( config , config_keys , item )

def get_updated_config( request ):
	config = DEFAULT_CONFIG
	if hasattr( request , "form" ) == False:
		return config

	update_config_with_form_key( request , "background_color" , config , [ "parser" , "our_background_textbox_hex" ] )
	update_config_with_form_key( request , "exported_width" , config , [ "parser" , "exported_width" ] )
	update_config_with_form_key( request , "exported_height" , config , [ "parser" , "exported_height" ] )
	update_config_with_form_key( request , "exported_image_dpi" , config , [ "parser" , "exported_image_dpi" ] )
	update_config_with_form_key( request , "image_scale_percentage" , config , [ "html" , "image_scale_percentage" ] )
	update_config_with_form_key( request , "unanswered_color" , config , [ "html" , "unanswered_color" ] )
	update_config_with_form_key( request , "answered_color" , config , [ "html" , "answered_color" ] )
	update_config_with_form_key( request , "text_color" , config , [ "html" , "text_color" ] )
	update_config_with_form_key( request , "text_font" , config , [ "html" , "text_font" ] )
	update_config_with_form_key( request , "text_x_offset_factor" , config , [ "html" , "text_x_offset_factor" ] )
	update_config_with_form_key( request , "text_y_offset_factor" , config , [ "html" , "text_y_offset_factor" ] )
	update_config_with_form_key( request , "base_hosted_url" , config , [ "html" , "base_hosted_url" ] )
	update_config_with_form_key( request , "interactive_typing_js" , config , [ "html" , "cdn" , "interactive_typing_js" ] )
	update_config_with_form_key( request , "interactive_drag_and_drop_js" , config , [ "html" , "cdn" , "interactive_drag_and_drop_js" ] )
	update_config_with_form_key( request , "jquery_ui_css" , config , [ "html" , "cdn" , "jquery_ui_css" , "url" ] )
	update_config_with_form_key( request , "jquery_ui_js" , config , [ "html" , "cdn" , "jquery_ui_js" , "url" ] )
	update_config_with_form_key( request , "jquery_js" , config , [ "html" , "cdn" , "jquery_js" , "url" ] )
	update_config_with_form_key( request , "bootstrap_css" , config , [ "html" , "cdn" , "bootstrap_css" , "url" ] )
	update_config_with_form_key( request , "bootstrap_bundle" , config , [ "html" , "cdn" , "bootstrap_bundle" , "url" ] )

	return config

@@app.route( "/host" , methods=[ "GET" ] ):
async def host( request: Request ):
	return sanic_html( f'''<!DOCTYPE html>
<!DOCTYPE html>
<html>
<head>
	<meta charset="utf-8">
	<meta name="viewport" content="width=device-width, initial-scale=1">
	<title>PowerPoint - Interactive Game Generator</title>
</head>
<body>
	<h1>PowerPoint Interactive Hosted Game Generator - Stage 1</h1>
	<h3>Instructions for Stage 1</h3>
	<ol>
		<li>Create a PowerPoint with TextBoxes that have all been filled with the same predetermined hex color</li>
		<li>Upload that .pptx file here</li>
		<li>Download Generated .pttx file that contain slides with the text removed</li>
		<li><a href="/host/stage/2">Go to Stage 2</a></li>
	</ol>
	<form enctype="multipart/form-data" action="/host/stage/1" method="POST">
		<input type="file" id="powerpoint" name="file"><br><br>
		<span>Textbox Background Color (Hex)</span>&nbsp&nbsp<input type="text" id="background_color" name="background_color" placeholder="0070C0"><br>
		<span>Exported Slide Image Width</span>&nbsp&nbsp<input type="text" id="exported_width" name="exported_width" placeholder="1920"><br>
		<span>Exported Slide Image Height</span>&nbsp&nbsp<input type="text" id="exported_height" name="exported_height" placeholder="1080"><br>
		<span>Exported Slide Image DPI</span>&nbsp&nbsp<input type="text" id="exported_image_dpi" name="exported_image_dpi" placeholder="144"><br>
		<span>Scale Percentage of Image</span>&nbsp&nbsp<input type="text" id="image_scale_percentage" name="image_scale_percentage" placeholder="60"><br>
		<span>Unanswered Color Outline</span>&nbsp&nbsp<input type="text" id="unanswered_color" name="unanswered_color" placeholder="red"><br>
		<span>Answered Color Outline</span>&nbsp&nbsp<input type="text" id="answered_color" name="answered_color" placeholder="#13E337"><br>
		<span>Typing Text Color</span>&nbsp&nbsp<input type="text" id="text_color" name="text_color" placeholder="white"><br>
		<span>Typing Text Font</span>&nbsp&nbsp<input type="text" id="text_font" name="text_font" placeholder="17px Arial"><br>
		<span>Typing Text X-Offset</span>&nbsp&nbsp<input type="text" id="text_x_offset_factor" name="text_x_offset_factor" placeholder="2"><br>
		<span>Typing Text Y-Offset</span>&nbsp&nbsp<input type="text" id="text_y_offset_factor" name="text_y_offset_factor" placeholder="3"><br>
		<span>Base URL of Hosted HTML</span>&nbsp&nbsp<input type="text" size="60" id="base_hosted_url" name="base_hosted_url" placeholder="https://39363.org/NOTES/WSU/2021/Fall/ANT3100/Interactive"><br>
		<span>CDN of Interactive Typing JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_typing_js" name="interactive_typing_js" placeholder="https://39363.org/CDN/NOTES/interactive_typing.js"><br>
		<span>CDN of Interactive DragAndDrop JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_drag_and_drop_js" name="interactive_drag_and_drop_js" placeholder="https://39363.org/CDN/NOTES/interactive_drag_and_drop.js"><br>
		<span>CDN of JQuery-UI CSS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_css" name="jquery_ui_css" placeholder="https://39363.org/CDN/jquery-ui.css"><br>
		<span>CDN of JQuery-UI JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_js" name="jquery_ui_js" placeholder="https://39363.org/CDN/jquery-ui.min.js"><br>
		<span>CDN of JQuery JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_js" name="jquery_js" placeholder="https://39363.org/CDN/jquery-3.6.0.min.js"><br>
		<span>CDN of Bootstrap CSS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_css" name="bootstrap_css" placeholder="https://39363.org/CDN/bootstrap.min.css"><br>
		<span>CDN of Bootstrap Bundle JS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_bundle" name="bootstrap_bundle" placeholder="https://39363.org/CDN/bootstrap.bundle.min.js"><br>
		<br>
		<input type="submit">
	</form>
</body>
</html>''')

@app.post( "/host/stage/1" , stream=True )
async def generate_host_stage_one( request ):
	try:
		start_time = time.time()
		result = bytes()
		print( "Uploading Host Stage 1 File" )
		while True:
			body = await request.stream.read()
			if body is None:
				break
			result += body
		# print( result )
		end_time = time.time()
		duration = round( end_time - start_time )
		duration_minutes = ( duration / 60 )
		durating_seconds = ( duration % 60 )
		print( request.stream )
		print( f"Upload Took {duration_minutes} minutes and {durating_seconds} seconds" )
		config = get_updated_config( request )
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_dir_posix = Path( temp_dir )
			input_file_path = temp_dir_posix.joinpath( "input.pptx" )
			output_powerpoint_path = temp_dir_posix.joinpath( f"Blank.pptx" )

			# 1.) Read Sent Bytes into .pptx file inside temp directory
			with open( str( input_file_path ) , "wb" ) as file:
				file.write( result )
			print( input_file_path )
			print( input_file_path.stat() )

			# 2.) Create Copy of PowerPoint With All text removed from boxes with correct fill color
			p = Presentation( input_file_path )
			p_clone = deepcopy( p )
			for slide_index , slide in enumerate( p_clone.slides ):
				for shape_index , shape in enumerate( slide.shapes ):
					if utils.validate_text_box( shape , config[ "parser" ][ "our_background_textbox_hex" ] ):
						print( f"{slide_index} === {shape_index} === valid text box" )
						shape.text_frame.text = ""
			p_clone.save( str( output_powerpoint_path ) )
			print( output_powerpoint_path )
			print( output_powerpoint_path.stat() )
			return await sanic_file(
				str( output_powerpoint_path ) ,
				mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation" ,
				filename=output_powerpoint_path.name
			)
			# Can't stream back because we are in temp file land , and it already deletes somehow idk
			# return await sanic_file_stream(
			# 	str( output_powerpoint_path ) ,
			# 	status=200 ,
			# 	chunk_size=4096 ,
			# 	mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation" ,
			# 	headers={ "Content-Length": str( input_file_path.stat().st_size ) } ,
			# 	filename=output_powerpoint_path.name ,
			# )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )


@app.route( "/local" , methods=[ "GET" ] )
async def local( request: Request ):
	return sanic_html( f'''<!DOCTYPE html>
<!DOCTYPE html>
<html>
<head>
	<meta charset="utf-8">
	<meta name="viewport" content="width=device-width, initial-scale=1">
	<title>PowerPoint - Interactive Game Generator</title>
</head>
<body>
	<h1>PowerPoint Interactive Local Game Generator - Stage 1</h1>
	<h3>Instructions for Stage 1</h3>
	<ol>
		<li>Create a PowerPoint with TextBoxes that have all been filled with the same predetermined hex color</li>
		<li>Upload that .pptx file here</li>
		<li>Download Generated .pttx file that contain slides with the text removed</li>
		<li><a href="/local/stage/2">Go to Stage 2</a></li>
	</ol>
	<form enctype="multipart/form-data" action="/local/stage/1" method="POST">
		<input type="file" id="powerpoint" name="file"><br><br>
		<span>Textbox Background Color (Hex)</span>&nbsp&nbsp<input type="text" id="background_color" name="background_color" placeholder="0070C0"><br>
		<span>Exported Slide Image Width</span>&nbsp&nbsp<input type="text" id="exported_width" name="exported_width" placeholder="1920"><br>
		<span>Exported Slide Image Height</span>&nbsp&nbsp<input type="text" id="exported_height" name="exported_height" placeholder="1080"><br>
		<span>Exported Slide Image DPI</span>&nbsp&nbsp<input type="text" id="exported_image_dpi" name="exported_image_dpi" placeholder="144"><br>
		<span>Scale Percentage of Image</span>&nbsp&nbsp<input type="text" id="image_scale_percentage" name="image_scale_percentage" placeholder="60"><br>
		<span>Unanswered Color Outline</span>&nbsp&nbsp<input type="text" id="unanswered_color" name="unanswered_color" placeholder="red"><br>
		<span>Answered Color Outline</span>&nbsp&nbsp<input type="text" id="answered_color" name="answered_color" placeholder="#13E337"><br>
		<span>Typing Text Color</span>&nbsp&nbsp<input type="text" id="text_color" name="text_color" placeholder="white"><br>
		<span>Typing Text Font</span>&nbsp&nbsp<input type="text" id="text_font" name="text_font" placeholder="17px Arial"><br>
		<span>Typing Text X-Offset</span>&nbsp&nbsp<input type="text" id="text_x_offset_factor" name="text_x_offset_factor" placeholder="2"><br>
		<span>Typing Text Y-Offset</span>&nbsp&nbsp<input type="text" id="text_y_offset_factor" name="text_y_offset_factor" placeholder="3"><br>
		<span>Base URL of Hosted HTML</span>&nbsp&nbsp<input type="text" size="60" id="base_hosted_url" name="base_hosted_url" placeholder="https://39363.org/NOTES/WSU/2021/Fall/ANT3100/Interactive"><br>
		<span>CDN of Interactive Typing JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_typing_js" name="interactive_typing_js" placeholder="https://39363.org/CDN/NOTES/interactive_typing.js"><br>
		<span>CDN of Interactive DragAndDrop JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_drag_and_drop_js" name="interactive_drag_and_drop_js" placeholder="https://39363.org/CDN/NOTES/interactive_drag_and_drop.js"><br>
		<span>CDN of JQuery-UI CSS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_css" name="jquery_ui_css" placeholder="https://39363.org/CDN/jquery-ui.css"><br>
		<span>CDN of JQuery-UI JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_js" name="jquery_ui_js" placeholder="https://39363.org/CDN/jquery-ui.min.js"><br>
		<span>CDN of JQuery JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_js" name="jquery_js" placeholder="https://39363.org/CDN/jquery-3.6.0.min.js"><br>
		<span>CDN of Bootstrap CSS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_css" name="bootstrap_css" placeholder="https://39363.org/CDN/bootstrap.min.css"><br>
		<span>CDN of Bootstrap Bundle JS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_bundle" name="bootstrap_bundle" placeholder="https://39363.org/CDN/bootstrap.bundle.min.js"><br>
		<br>
		<input type="submit">
	</form>
</body>
</html>''')

# https://stackoverflow.com/questions/4212861/what-is-a-correct-mime-type-for-docx-pptx-etc
@app.post( "/local/stage/1" , stream=True )
async def generate_local_stage_one( request ):
	try:
		start_time = time.time()
		result = bytes()
		print( "Uploading Local Stage 1 File" )
		while True:
			body = await request.stream.read()
			if body is None:
				break
			result += body
		# print( result )
		end_time = time.time()
		duration = round( end_time - start_time )
		duration_minutes = ( duration / 60 )
		durating_seconds = ( duration % 60 )
		print( request.stream )
		print( f"Upload Took {duration_minutes} minutes and {durating_seconds} seconds" )
		config = get_updated_config( request )
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_dir_posix = Path( temp_dir )
			input_file_path = temp_dir_posix.joinpath( "input.pptx" )
			output_powerpoint_path = temp_dir_posix.joinpath( f"Blank.pptx" )

			# 1.) Read Sent Bytes into .pptx file inside temp directory
			with open( str( input_file_path ) , "wb" ) as file:
				file.write( result )
			print( input_file_path )
			print( input_file_path.stat() )

			# 2.) Create Copy of PowerPoint With All text removed from boxes with correct fill color
			p = Presentation( input_file_path )
			p_clone = deepcopy( p )
			for slide_index , slide in enumerate( p_clone.slides ):
				for shape_index , shape in enumerate( slide.shapes ):
					if utils.validate_text_box( shape , config[ "parser" ][ "our_background_textbox_hex" ] ):
						print( f"{slide_index} === {shape_index} === valid text box" )
						shape.text_frame.text = ""
			p_clone.save( str( output_powerpoint_path ) )
			print( output_powerpoint_path )
			print( output_powerpoint_path.stat() )
			return await sanic_file(
				str( output_powerpoint_path ) ,
				mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation" ,
				filename=output_powerpoint_path.name
			)
			# Can't stream back because we are in temp file land , and it already deletes somehow idk
			# return await sanic_file_stream(
			# 	str( output_powerpoint_path ) ,
			# 	status=200 ,
			# 	chunk_size=4096 ,
			# 	mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation" ,
			# 	headers={ "Content-Length": str( input_file_path.stat().st_size ) } ,
			# 	filename=output_powerpoint_path.name ,
			# )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )


@app.route( "/local/stage/2" , methods=[ "GET" ] )
async def upload( request: Request ):
	return sanic_html( f'''<!DOCTYPE html>
<!DOCTYPE html>
<html>
<head>
	<meta charset="utf-8">
	<meta name="viewport" content="width=device-width, initial-scale=1">
	<title>PowerPoint - Interactive Game Generator</title>
</head>
<body>
	<h1>PowerPoint Interactive Local Game Generator - Stage 2</h1>
	<h3>Instructions for Stage 2 </h3>
	<ol>
		<li>Open the -Blank.pptx that you should of just downloaded in PowerPoint</li>
		<li>File --> Export --> File Format: JPEG , Save Every Slide , Width: 1920 , Height: 1080</li>
		<li>Create .zip archive of the orginal .pptx you uploaded in Stage 1 and the folder of jpegs it just generated</li>
		<li>Upload that .zip file here</a></li>
	</ol>
	<form enctype="multipart/form-data" action="/local/stage/2" method="POST">
		<input type="file" id="powerpoint" name="file"><br><br>
		<span>Textbox Background Color (Hex)</span>&nbsp&nbsp<input type="text" id="background_color" name="background_color" placeholder="0070C0"><br>
		<span>Exported Slide Image Width</span>&nbsp&nbsp<input type="text" id="exported_width" name="exported_width" placeholder="1920"><br>
		<span>Exported Slide Image Height</span>&nbsp&nbsp<input type="text" id="exported_height" name="exported_height" placeholder="1080"><br>
		<span>Exported Slide Image DPI</span>&nbsp&nbsp<input type="text" id="exported_image_dpi" name="exported_image_dpi" placeholder="144"><br>
		<span>Scale Percentage of Image</span>&nbsp&nbsp<input type="text" id="image_scale_percentage" name="image_scale_percentage" placeholder="60"><br>
		<span>Unanswered Color Outline</span>&nbsp&nbsp<input type="text" id="unanswered_color" name="unanswered_color" placeholder="red"><br>
		<span>Answered Color Outline</span>&nbsp&nbsp<input type="text" id="answered_color" name="answered_color" placeholder="#13E337"><br>
		<span>Typing Text Color</span>&nbsp&nbsp<input type="text" id="text_color" name="text_color" placeholder="white"><br>
		<span>Typing Text Font</span>&nbsp&nbsp<input type="text" id="text_font" name="text_font" placeholder="17px Arial"><br>
		<span>Typing Text X-Offset</span>&nbsp&nbsp<input type="text" id="text_x_offset_factor" name="text_x_offset_factor" placeholder="2"><br>
		<span>Typing Text Y-Offset</span>&nbsp&nbsp<input type="text" id="text_y_offset_factor" name="text_y_offset_factor" placeholder="3"><br>
		<span>Base URL of Hosted HTML</span>&nbsp&nbsp<input type="text" size="60" id="base_hosted_url" name="base_hosted_url" placeholder="https://39363.org/NOTES/WSU/2021/Fall/ANT3100/Interactive"><br>
		<span>CDN of Interactive Typing JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_typing_js" name="interactive_typing_js" placeholder="https://39363.org/CDN/NOTES/interactive_typing.js"><br>
		<span>CDN of Interactive DragAndDrop JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_drag_and_drop_js" name="interactive_drag_and_drop_js" placeholder="https://39363.org/CDN/NOTES/interactive_drag_and_drop.js"><br>
		<span>CDN of JQuery-UI CSS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_css" name="jquery_ui_css" placeholder="https://39363.org/CDN/jquery-ui.css"><br>
		<span>CDN of JQuery-UI JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_js" name="jquery_ui_js" placeholder="https://39363.org/CDN/jquery-ui.min.js"><br>
		<span>CDN of JQuery JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_js" name="jquery_js" placeholder="https://39363.org/CDN/jquery-3.6.0.min.js"><br>
		<span>CDN of Bootstrap CSS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_css" name="bootstrap_css" placeholder="https://39363.org/CDN/bootstrap.min.css"><br>
		<span>CDN of Bootstrap Bundle JS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_bundle" name="bootstrap_bundle" placeholder="https://39363.org/CDN/bootstrap.bundle.min.js"><br>
		<br>
		<input type="submit">
	</form>
</body>
</html>''')

@app.post( "/local/stage/2" , stream=True )
async def generate_local_stage_two( request: Request ):
	try:
		start_time = time.time()
		result = bytes()
		print( "Uploading Local Stage 2 File" )
		while True:
			body = await request.stream.read()
			if body is None:
				break
			result += body
		# print( result )
		end_time = time.time()
		duration = round( end_time - start_time )
		duration_minutes = ( duration / 60 )
		durating_seconds = ( duration % 60 )
		print( request.stream )
		print( f"Upload Took {duration_minutes} minutes and {durating_seconds} seconds" )

		config = get_updated_config( request )

		with tempfile.TemporaryDirectory() as temp_dir:

			temp_dir_posix = Path( temp_dir )
			input_zip_file_path = temp_dir_posix.joinpath( "input.zip" )
			# output_powerpoint_path = temp_dir_posix.joinpath( f"{input_file_name_stem}-Blank.pptx" )

			# 1.) Read Sent Bytes into .zip file inside temp directory
			with open( str( input_zip_file_path ) , "wb"  ) as file:
				file.write( result )
			print( input_zip_file_path )
			print( input_zip_file_path.stat() )

			# 2.) Extract the Zip
			with zipfile.ZipFile( str( input_zip_file_path ) , "r" ) as zip_ref:
				zip_ref.extractall( str( temp_dir_posix ) )

			# 3.) Sort Files and Folders Contained in Zip
			files_and_folders = temp_dir_posix.glob( '*' )
			uploaded_powerpoint_path = False
			image_paths = []
			for x in temp_dir_posix.iterdir():
				if x.is_file():
					if x.suffix == ".pptx":
						uploaded_powerpoint_path = x
				elif x.is_dir():
					if x.stem == "__MACOSX":
						continue
					jpegs = x.glob( "*" )
					jpegs = [ i for i in jpegs if i.is_file() ]
					for index , i in enumerate( jpegs ):
						if i.suffix == ".jpg":
							image_paths.append( x.joinpath( f"Slide{index+1}.jpg" ) )
						elif i.suffix == ".jpeg":
							image_paths.append( x.joinpath( f"Slide{index+1}.jpeg" ) )
			print( uploaded_powerpoint_path )
			print( image_paths )


			# 4.) Create Placeholder Output Directories
			# generated_output_base_path = temp_dir_posix.joinpath( f"{input_file_name_stem} - Interactive" )
			generated_output_base_path = temp_dir_posix.joinpath( f"Interactive" )
			generated_output_base_path.mkdir( parents=True , exist_ok=True )
			output_images_path = generated_output_base_path.joinpath( "images" )
			output_images_path.mkdir( parents=True , exist_ok=True )
			output_html_path = generated_output_base_path.joinpath( "html" )
			output_html_path.mkdir( parents=True , exist_ok=True )
			print( output_images_path )
			for index , image_path in enumerate( image_paths ):
				shutil.copyfile( image_path , output_images_path.joinpath( image_path.name ) )
			output_js_path = generated_output_base_path.joinpath( "js" )
			shutil.copytree( "./js" , output_js_path )
			output_js_path.mkdir( parents=True , exist_ok=True )
			output_css_path = generated_output_base_path.joinpath( "css" )
			shutil.copytree( "./css" , output_css_path )
			output_css_path.mkdir( parents=True , exist_ok=True )

			image_maps_in_slides = compute_image_maps.compute( str( uploaded_powerpoint_path ) , config[ "parser" ] )
			print( len( image_maps_in_slides ) )

			if len( image_paths ) != len( image_maps_in_slides ):
				print( "Lengths of Slide Images and Slide Image Maps Don't Match" )
				# print( "This could mean some slides don't have any thing but" )
				sys.exit( 1 )

			print( "here 2" )
			image_objects = []
			for slide_index , image_maps_in_slide in enumerate( image_maps_in_slides ):
				print( f"\nSlide === {slide_index+1}" )
				image_objects.append([
					f"Slide" ,
					f"../../images/{image_paths[ slide_index ].name}" ,
					image_maps_in_slide
				])
			print( "here 3" )
			pprint( image_objects )

			# 3.) Generate Interactive Games
			config[ "html" ][ "base_hosted_url" ] = "./"
			config[ "html" ][ "cdn" ][ "jquery_ui_css"][ "url" ] = f"../../css/jquery-ui.css"
			config[ "html" ][ "cdn" ][ "jquery_ui_js"][ "url" ] = f"../../js/jquery-ui.min.js"
			config[ "html" ][ "cdn" ][ "jquery_js"][ "url" ] = f"../../js/jquery-3.6.0.min.js"
			config[ "html" ][ "cdn" ][ "bootstrap_css"][ "url" ] = f"../../css/bootstrap.min.css"
			config[ "html" ][ "cdn" ][ "bootstrap_bundle"][ "url" ] = f"../../js/bootstrap.bundle.min.js"
			config[ "html" ][ "cdn" ][ "interactive_typing_js" ] = f"../../js/interactive_typing.js"
			config[ "html" ][ "cdn" ][ "interactive_drag_and_drop_js" ] = f"../../js/interactive_drag_and_drop.js"
			html_options = config[ "html" ]
			html_options[ "title" ] = uploaded_powerpoint_path.stem.replace( " " , "-" )
			html_options[ "images" ] = image_objects
			html_options[ "output_base_dir" ] = output_html_path
			interactive_notes_generator.generate( html_options , False )

			# 4.) Zip Everything and Return
			shutil.make_archive( temp_dir_posix.joinpath( f"{generated_output_base_path.stem}" ) , "zip" , str( generated_output_base_path ) )
			output_zip_path = temp_dir_posix.joinpath( f"{generated_output_base_path.stem}.zip" )
			return await sanic_file(
				str( output_zip_path ) ,
				mime_type="application/zip" ,
				filename=f"{generated_output_base_path.stem}.zip"
			)
			# return sanic_json( dict( ok="ok" ) , status=200 )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )


if __name__ == "__main__":
	app.run( host="0.0.0.0" , port="9379" )