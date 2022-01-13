#!/usr/bin/env python3
import sys
import uuid
import time
import json
import base64
from pathlib import Path
from pprint import pprint
import tempfile
import zipfile
import shutil
import requests
from binascii import b2a_hex

from sanic import Sanic
from sanic import Blueprint
from sanic.response import html as sanic_html
from sanic.response import raw as sanic_raw
from sanic.response import json as sanic_json
from sanic.response import file as sanic_file
from sanic.response import file_stream as sanic_file_stream
from sanic.response import stream as sanic_stream
from sanic.request import Request

from sanic_jwt_extended import JWT , jwt_required
from sanic_jwt_extended.tokens import Token

import utils
import compute_image_maps
import image_uploader
import interactive_notes_generator

from pptx import Presentation
from copy import deepcopy

import pytz
import datetime
from datetime import timedelta
from dateutil.relativedelta import relativedelta
import magic
from ulid import ULID
import jwt


time_zone = pytz.timezone( "US/Eastern" )
app = Sanic( __name__ )
# https://pypi.org/project/python-magic/
# sudo apt-get install libmagic1
# brew install libmagic
# pip install python-magic-bin
mime = magic.Magic( mime=True )
DEFAULT_CONFIG = utils.read_json( sys.argv[1] )

app.static( "/test/host/static/js" , "./js" )
app.static( "/test/host/static/css" , "./css" )

# app.static( "/css" , "./css" )

# https://sanic-jwt-extended.seonghyeon.dev/config_options.html
# https://pyjwt.readthedocs.io/en/latest/algorithms.html
with JWT.initialize( app ) as manager:
	manager.config.secret_key = DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ]
	manager.config.public_key = DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_public_key" ]
	manager.config.private_key = DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_private_key" ]
	manager.config.algorithm = "HS512"
	manager.config.jwt_header_key = "token"
	manager.config.jwt_header_prefix = "Bearer"
	manager.config.jwt_cookie = "ppt_interactive_image_upload_server"
	manager.config.access_token_expires = timedelta( weeks=( 52 * 10 ) )

def update_config_with_form_key( request , form_key , config , *config_keys ):
	item = request.form.get( form_key )
	if item != None and len( item ) > 0:
		utils.set_nested_dictionary_value( config , config_keys , item )
def get_updated_config( request ):
	config = DEFAULT_CONFIG
	if hasattr( request , "form" ) == False:
		return config
	update_config_with_form_key( request , "background_color" , config , [ "parser" , "our_background_textbox_hex" ] )
	update_config_with_form_key( request , "exported_width" , config , [ "parser" , "exported_width" ] )
	update_config_with_form_key( request , "exported_height" , config , [ "parser" , "exported_height" ] )
	update_config_with_form_key( request , "exported_image_dpi" , config , [ "parser" , "exported_image_dpi" ] )
	update_config_with_form_key( request , "image_scale_percentage" , config , [ "html" , "image_scale_percentage" ] )
	update_config_with_form_key( request , "unanswered_color" , config , [ "html" , "unanswered_color" ] )
	update_config_with_form_key( request , "answered_color" , config , [ "html" , "answered_color" ] )
	update_config_with_form_key( request , "text_color" , config , [ "html" , "text_color" ] )
	update_config_with_form_key( request , "text_font" , config , [ "html" , "text_font" ] )
	update_config_with_form_key( request , "text_x_offset_factor" , config , [ "html" , "text_x_offset_factor" ] )
	update_config_with_form_key( request , "text_y_offset_factor" , config , [ "html" , "text_y_offset_factor" ] )
	update_config_with_form_key( request , "base_hosted_url" , config , [ "html" , "base_hosted_url" ] )
	update_config_with_form_key( request , "interactive_typing_js" , config , [ "html" , "cdn" , "interactive_typing_js" ] )
	update_config_with_form_key( request , "interactive_drag_and_drop_js" , config , [ "html" , "cdn" , "interactive_drag_and_drop_js" ] )
	update_config_with_form_key( request , "jquery_ui_css" , config , [ "html" , "cdn" , "jquery_ui_css" , "url" ] )
	update_config_with_form_key( request , "jquery_ui_js" , config , [ "html" , "cdn" , "jquery_ui_js" , "url" ] )
	update_config_with_form_key( request , "jquery_js" , config , [ "html" , "cdn" , "jquery_js" , "url" ] )
	update_config_with_form_key( request , "bootstrap_css" , config , [ "html" , "cdn" , "bootstrap_css" , "url" ] )
	update_config_with_form_key( request , "bootstrap_bundle" , config , [ "html" , "cdn" , "bootstrap_bundle" , "url" ] )

	## Addon
	# secret_key = request.form.get( "secret_key" )
	# if item != None and len( item ) > 0:
	# 	config[ "secret_key" ] = secret_key
	# just generate new secret key every time
	return config

@app.route( "/" , methods=[ "GET" ] )
async def local( request: Request ):
	return sanic_html( f'''<!DOCTYPE html>
<!DOCTYPE html>
<html>
<head>
	<meta charset="utf-8">
	<meta name="viewport" content="width=device-width, initial-scale=1">
	<title>Test Uploader</title>
</head>
<body>
	<h1>Test Uploader</h1>
	<form enctype="multipart/form-data" action="/test/host/stage/1" method="POST">
		<input type="file" id="powerpoint" name="file"><br><br>
		<input type="submit">
	</form>
</body>
</html>''')
@app.post( "/test/upload" , stream=True )
async def test_upload( request ):
	try:

		# 1.) Injest octet-stream
		start_time = datetime.datetime.now().astimezone( time_zone )
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_dir_posix = Path( temp_dir )
			input_file_path = temp_dir_posix.joinpath( "input" )
			print( "Uploading Test File" )
			first_128 = bytes()
			with open( str( input_file_path ) , "wb" ) as file:
				while True:
					body = await request.stream.read()
					if body is None:
						break
					if len( first_128 ) < 129:
						first_128 += body
					file.write( body )
			end_time = datetime.datetime.now().astimezone( time_zone )
			duration = relativedelta( end_time , start_time )
			duration_milliseconds = round( duration.microseconds / 1000 )
			print( "The Upload Took %d year %d month %d days %d hours %d minutes %d seconds %d milliseconds" % (
				duration.years , duration.months , duration.days , duration.hours , duration.minutes ,
				duration.seconds , duration_milliseconds
			))

			# 2.) Discover File Type
			print( input_file_path )
			file_name = False
			file_extension = False
			file_content_type = False
			try:
				test = first_128[ 0 : 500 ].decode( "utf-8" , "ignore" )
				file_name = test.split( "filename=" )[ 1 ].split( "\n" )[ 0 ][ 1:-2 ]
				file_extension = file_name.split( "." )[ -1 ]
				file_content_type = test.split( "Content-Type: " )[ 1 ].split( "\n" )[ 0 ]
			except Exception as file_type_decode_exception:
				print( file_type_decode_exception )
			print( file_name , file_extension , file_content_type )
			# https://python-forum.io/thread-11941-post-54130.html#pid54130
			# if file_extension != False:
			# 	print( "????" )
			# 	input_file_path.rename( temp_dir_posix.joinpath( file_name ) )
			# 	print( input_file_path )
			print( input_file_path.stat() )

			return sanic_json( dict( size=str( input_file_path.stat() ) ) , status=200 )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )


# def determine_image_type( self , stream_first_4_bytes ):
# 	file_type = None
# 	bytes_as_hex = b2a_hex( stream_first_4_bytes ).decode()
# 	if bytes_as_hex.startswith( 'ffd8' ):
# 		file_type = '.jpeg'
# 	elif bytes_as_hex == '89504e47':
# 		file_type = '.png'
# 	elif bytes_as_hex == '47494638':
# 		file_type = '.gif'
# 	elif bytes_as_hex.startswith('424d'):
# 		file_type = '.bmp'
# 	return file_type

@app.route( "/test/host" , methods=[ "GET" ] )
async def local( request: Request ):
	return sanic_html( f'''<!DOCTYPE html>
<!DOCTYPE html>
<html>
<head>
	<meta charset="utf-8">
	<meta name="viewport" content="width=device-width, initial-scale=1">
	<title>PowerPoint - Interactive Game Generator</title>
</head>
<body>
	<h1>PowerPoint Interactive Hosted Game Generator - Stage 1</h1>
	<h3>Instructions for Stage 1</h3>
	<ol>
		<li>Create a PowerPoint with TextBoxes that have all been filled with the same predetermined hex color</li>
		<li>Upload that .pptx file here</li>
		<li>Download Generated .pttx file that contain slides with the text removed</li>
		<li><a href="/test/host/stage/2">Go to Stage 2</a></li>
	</ol>
	<form enctype="multipart/form-data" action="/test/host/stage/1" method="POST">
		<input type="file" id="powerpoint" name="file"><br><br>
		<span>Textbox Background Color (Hex)</span>&nbsp&nbsp<input type="text" id="background_color" name="background_color" placeholder="0070C0"><br>
		<span>Exported Slide Image Width</span>&nbsp&nbsp<input type="text" id="exported_width" name="exported_width" placeholder="1920"><br>
		<span>Exported Slide Image Height</span>&nbsp&nbsp<input type="text" id="exported_height" name="exported_height" placeholder="1080"><br>
		<span>Exported Slide Image DPI</span>&nbsp&nbsp<input type="text" id="exported_image_dpi" name="exported_image_dpi" placeholder="144"><br>
		<span>Scale Percentage of Image</span>&nbsp&nbsp<input type="text" id="image_scale_percentage" name="image_scale_percentage" placeholder="60"><br>
		<span>Unanswered Color Outline</span>&nbsp&nbsp<input type="text" id="unanswered_color" name="unanswered_color" placeholder="red"><br>
		<span>Answered Color Outline</span>&nbsp&nbsp<input type="text" id="answered_color" name="answered_color" placeholder="#13E337"><br>
		<span>Typing Text Color</span>&nbsp&nbsp<input type="text" id="text_color" name="text_color" placeholder="white"><br>
		<span>Typing Text Font</span>&nbsp&nbsp<input type="text" id="text_font" name="text_font" placeholder="17px Arial"><br>
		<span>Typing Text X-Offset</span>&nbsp&nbsp<input type="text" id="text_x_offset_factor" name="text_x_offset_factor" placeholder="2"><br>
		<span>Typing Text Y-Offset</span>&nbsp&nbsp<input type="text" id="text_y_offset_factor" name="text_y_offset_factor" placeholder="3"><br>
		<span>Base URL of Hosted HTML</span>&nbsp&nbsp<input type="text" size="60" id="base_hosted_url" name="base_hosted_url" placeholder="https://39363.org/NOTES/WSU/2021/Fall/ANT3100/Interactive"><br>
		<span>CDN of Interactive Typing JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_typing_js" name="interactive_typing_js" placeholder="https://39363.org/CDN/NOTES/interactive_typing.js"><br>
		<span>CDN of Interactive DragAndDrop JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_drag_and_drop_js" name="interactive_drag_and_drop_js" placeholder="https://39363.org/CDN/NOTES/interactive_drag_and_drop.js"><br>
		<span>CDN of JQuery-UI CSS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_css" name="jquery_ui_css" placeholder="https://39363.org/CDN/jquery-ui.css"><br>
		<span>CDN of JQuery-UI JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_js" name="jquery_ui_js" placeholder="https://39363.org/CDN/jquery-ui.min.js"><br>
		<span>CDN of JQuery JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_js" name="jquery_js" placeholder="https://39363.org/CDN/jquery-3.6.0.min.js"><br>
		<span>CDN of Bootstrap CSS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_css" name="bootstrap_css" placeholder="https://39363.org/CDN/bootstrap.min.css"><br>
		<span>CDN of Bootstrap Bundle JS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_bundle" name="bootstrap_bundle" placeholder="https://39363.org/CDN/bootstrap.bundle.min.js"><br>
		<br>
		<input type="submit">
	</form>
</body>
</html>''')
@app.post( "/test/host/stage/1" , stream=True )
async def test_upload_stage_1( request ):
	try:

		config = get_updated_config( request )

		# 1.) Injest octet-stream
		start_time = datetime.datetime.now().astimezone( time_zone )
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_dir_posix = Path( temp_dir )
			input_file_path = temp_dir_posix.joinpath( "input" )
			print( "Uploading Test File" )
			first_128 = bytes()
			with open( str( input_file_path ) , "wb" ) as file:
				while True:
					body = await request.stream.read()
					if body is None:
						break
					if len( first_128 ) < 129:
						first_128 += body
					file.write( body )
			end_time = datetime.datetime.now().astimezone( time_zone )
			duration = relativedelta( end_time , start_time )
			duration_milliseconds = round( duration.microseconds / 1000 )
			print( "The Upload Took %d year %d month %d days %d hours %d minutes %d seconds %d milliseconds" % (
				duration.years , duration.months , duration.days , duration.hours , duration.minutes ,
				duration.seconds , duration_milliseconds
			))

			# 2.) Discover File Type
			print( input_file_path )
			file_name = False
			file_extension = False
			file_name_stem = False
			file_content_type = False
			try:
				test = first_128[ 0 : 500 ].decode( "utf-8" , "ignore" )
				file_name = test.split( "filename=" )[ 1 ].split( "\n" )[ 0 ][ 1:-2 ]
				file_extension = file_name.split( "." )[ -1 ]
				file_name_stem = ".".join( file_name.split( "." )[ 0 : -1 ] )
				file_content_type = test.split( "Content-Type: " )[ 1 ].split( "\n" )[ 0 ]
			except Exception as file_type_decode_exception:
				print( file_type_decode_exception )
			print( file_name , file_extension , file_content_type )
			print( input_file_path.stat() )
			if file_name != False:
				output_powerpoint_path = temp_dir_posix.joinpath( f"{file_name_stem}-Blank.pptx" )
			else:
				output_powerpoint_path = temp_dir_posix.joinpath( "Blank.pptx" )

			# 3.) Create Copy of PowerPoint With All text removed from boxes with correct fill color
			p = Presentation( input_file_path )
			p_clone = deepcopy( p )
			for slide_index , slide in enumerate( p_clone.slides ):
				for shape_index , shape in enumerate( slide.shapes ):
					if utils.validate_text_box( shape , config[ "parser" ][ "our_background_textbox_hex" ] ):
						print( f"{slide_index} === {shape_index} === valid text box" )
						shape.text_frame.text = ""
			p_clone.save( str( output_powerpoint_path ) )
			print( output_powerpoint_path )
			print( output_powerpoint_path.stat() )
			return await sanic_file(
				str( output_powerpoint_path ) ,
				mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation" ,
				filename=output_powerpoint_path.name
			)
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )


@app.route( "/test/host/stage/2" , methods=[ "GET" ] )
async def test_upload_stage_2( request: Request ):
	return sanic_html( f'''<!DOCTYPE html>
<!DOCTYPE html>
<html>
<head>
	<meta charset="utf-8">
	<meta name="viewport" content="width=device-width, initial-scale=1">
	<title>PowerPoint - Interactive Game Generator</title>
</head>
<body>
	<h1>PowerPoint Interactive Hosted Game Generator - Stage 2</h1>
	<h3>Instructions for Stage 2 </h3>
	<ol>
		<li>Open the -Blank.pptx that you should of just downloaded in PowerPoint</li>
		<li>File --> Export --> File Format: JPEG , Save Every Slide , Width: 1920 , Height: 1080</li>
		<li>Create .zip archive of the orginal .pptx you uploaded in Stage 1 and the folder of jpegs it just generated</li>
		<li>Upload that .zip file here</a></li>
	</ol>
	<form enctype="multipart/form-data" action="/test/host/stage/2" method="POST">
		<input type="file" id="powerpoint" name="file"><br><br>
		<span>Textbox Background Color (Hex)</span>&nbsp&nbsp<input type="text" id="background_color" name="background_color" placeholder="0070C0"><br>
		<span>Exported Slide Image Width</span>&nbsp&nbsp<input type="text" id="exported_width" name="exported_width" placeholder="1920"><br>
		<span>Exported Slide Image Height</span>&nbsp&nbsp<input type="text" id="exported_height" name="exported_height" placeholder="1080"><br>
		<span>Exported Slide Image DPI</span>&nbsp&nbsp<input type="text" id="exported_image_dpi" name="exported_image_dpi" placeholder="144"><br>
		<span>Scale Percentage of Image</span>&nbsp&nbsp<input type="text" id="image_scale_percentage" name="image_scale_percentage" placeholder="60"><br>
		<span>Unanswered Color Outline</span>&nbsp&nbsp<input type="text" id="unanswered_color" name="unanswered_color" placeholder="red"><br>
		<span>Answered Color Outline</span>&nbsp&nbsp<input type="text" id="answered_color" name="answered_color" placeholder="#13E337"><br>
		<span>Typing Text Color</span>&nbsp&nbsp<input type="text" id="text_color" name="text_color" placeholder="white"><br>
		<span>Typing Text Font</span>&nbsp&nbsp<input type="text" id="text_font" name="text_font" placeholder="17px Arial"><br>
		<span>Typing Text X-Offset</span>&nbsp&nbsp<input type="text" id="text_x_offset_factor" name="text_x_offset_factor" placeholder="2"><br>
		<span>Typing Text Y-Offset</span>&nbsp&nbsp<input type="text" id="text_y_offset_factor" name="text_y_offset_factor" placeholder="3"><br>
		<span>Base URL of Hosted HTML</span>&nbsp&nbsp<input type="text" size="60" id="base_hosted_url" name="base_hosted_url" placeholder="https://39363.org/NOTES/WSU/2021/Fall/ANT3100/Interactive"><br>
		<span>CDN of Interactive Typing JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_typing_js" name="interactive_typing_js" placeholder="https://39363.org/CDN/NOTES/interactive_typing.js"><br>
		<span>CDN of Interactive DragAndDrop JS</span>&nbsp&nbsp<input type="text" size="60" id="interactive_drag_and_drop_js" name="interactive_drag_and_drop_js" placeholder="https://39363.org/CDN/NOTES/interactive_drag_and_drop.js"><br>
		<span>CDN of JQuery-UI CSS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_css" name="jquery_ui_css" placeholder="https://39363.org/CDN/jquery-ui.css"><br>
		<span>CDN of JQuery-UI JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_ui_js" name="jquery_ui_js" placeholder="https://39363.org/CDN/jquery-ui.min.js"><br>
		<span>CDN of JQuery JS</span>&nbsp&nbsp<input type="text" size="60" id="jquery_js" name="jquery_js" placeholder="https://39363.org/CDN/jquery-3.6.0.min.js"><br>
		<span>CDN of Bootstrap CSS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_css" name="bootstrap_css" placeholder="https://39363.org/CDN/bootstrap.min.css"><br>
		<span>CDN of Bootstrap Bundle JS</span>&nbsp&nbsp<input type="text" size="60" id="bootstrap_bundle" name="bootstrap_bundle" placeholder="https://39363.org/CDN/bootstrap.bundle.min.js"><br>
		<br>
		<input type="submit">
	</form>
</body>
</html>''')
@app.post( "/test/host/stage/2" , stream=True )
async def test_upload_stage_2( request ):
	try:

		config = get_updated_config( request )

		# 1.) Injest octet-stream
		start_time = datetime.datetime.now().astimezone( time_zone )
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_dir_posix = Path( temp_dir )
			input_file_path = temp_dir_posix.joinpath( "input.zip" )
			print( "Uploading Test File" )
			first_128 = bytes()
			with open( str( input_file_path ) , "wb" ) as file:
				while True:
					body = await request.stream.read()
					if body is None:
						break
					if len( first_128 ) < 129:
						first_128 += body
					file.write( body )
			end_time = datetime.datetime.now().astimezone( time_zone )
			duration = relativedelta( end_time , start_time )
			duration_milliseconds = round( duration.microseconds / 1000 )
			print( "The Upload Took %d year %d month %d days %d hours %d minutes %d seconds %d milliseconds" % (
				duration.years , duration.months , duration.days , duration.hours , duration.minutes ,
				duration.seconds , duration_milliseconds
			))

			# 2.) Discover File Type
			print( input_file_path )
			file_name = False
			file_extension = False
			file_name_stem = False
			file_content_type = False
			try:
				test = first_128[ 0 : 500 ].decode( "utf-8" , "ignore" )
				file_name = test.split( "filename=" )[ 1 ].split( "\n" )[ 0 ][ 1:-2 ]
				file_extension = file_name.split( "." )[ -1 ]
				file_name_stem = ".".join( file_name.split( "." )[ 0 : -1 ] )
				file_content_type = test.split( "Content-Type: " )[ 1 ].split( "\n" )[ 0 ]
			except Exception as file_type_decode_exception:
				print( file_type_decode_exception )
			print( file_name , file_extension , file_content_type )
			print( input_file_path.stat() )
			# if file_name != False:
			# 	output_powerpoint_path = temp_dir_posix.joinpath( f"{file_name_stem}-Blank.pptx" )
			# else:
			# 	output_powerpoint_path = temp_dir_posix.joinpath( "Blank.pptx" )

			# 3.) Extract the Zip
			with zipfile.ZipFile( str( input_file_path ) , "r" ) as zip_ref:
				zip_ref.extractall( str( temp_dir_posix ) )

			# 4.) Sort Files and Folders Contained in Zip
			files_and_folders = temp_dir_posix.glob( '*' )
			uploaded_powerpoint_path = False
			image_paths = []
			for x in temp_dir_posix.iterdir():
				if x.is_file():
					if x.suffix == ".pptx":
						uploaded_powerpoint_path = x
				elif x.is_dir():
					if x.stem == "__MACOSX":
						continue
					jpegs = x.glob( "*" )
					jpegs = [ i for i in jpegs if i.is_file() ]
					for index , i in enumerate( jpegs ):
						if i.suffix == ".jpg":
							image_paths.append( x.joinpath( f"Slide{index+1}.jpg" ) )
						elif i.suffix == ".jpeg":
							image_paths.append( x.joinpath( f"Slide{index+1}.jpeg" ) )
			print( uploaded_powerpoint_path )
			print( image_paths )

			## Addon , generate secret box key unique to this .pptx
			SECRET_BOX_KEY = utils.secret_box_generate_new_key()
			print( SECRET_BOX_KEY )

			# 5.) Copy Uploaded Images to Hosted Storage Directory
			output_images_path = Path( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "local_image_storage_path" ] )
			output_blob_path = Path( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "local_blob_storage_path" ] )
			image_ulids = []
			for index , image_path in enumerate( image_paths ):
				#shutil.copyfile( image_path , output_images_path.joinpath( image_path.name ) )
				## Addon , encrypt images
				with open( str( image_path ) , "rb" ) as img_file:
					image_bytes = img_file.read()
					image_bytes_b64_bytes = base64.b64encode( image_bytes )
					image_bytes_b64_string = image_bytes_b64_bytes.decode()
					image_ulid = str( ULID() )
					image_ulids.append( image_ulid )
					print( type( image_bytes_b64_string ) )
					utils.write_json( str( output_images_path.joinpath( f"{image_ulid}.json" ) ) , {
						"sealed": utils.secret_box_seal( SECRET_BOX_KEY , image_bytes_b64_string )
					})
					print( "3" )
			print( "5" )

			# 6.) Compute Image Maps
			image_maps_in_slides = compute_image_maps.compute( str( uploaded_powerpoint_path ) , config[ "parser" ] )
			print( len( image_maps_in_slides ) )
			if len( image_paths ) != len( image_maps_in_slides ):
				print( "Lengths of Slide Images and Slide Image Maps Don't Match" )
				# print( "This could mean some slides don't have any thing but" )
				sys.exit( 1 )


			## Addon.) Create "blob" object that can just be passed either to DragAndDrop or Typing HTML file
			## Contains image-maps , encrypted hosted-image-file-path objects , ALWAYS maintain expansion opportunity , add meta
			## Addon.) Generate ULID for each ImagePath
			## 			 Then , seal ULID in a JWT Token , this will be appended as a url parameter ?p="asdf"
			# https://pyjwt.readthedocs.io/en/latest/usage.html
			# https://pyjwt.readthedocs.io/en/latest/algorithms.html?highlight=algorithm
			# Ran into problems with RS512 ,
			# you can generate a keypair : openssl ecparam -name secp521r1 -genkey -noout -out private.ec.key
			# and then store it base64 encoded in config.json
			# , but on mac osx where we are testing : '_EllipticCurvePrivateKey' object has no attribute 'verify'
			# so changed back to HS512
			print( "here-1" )
			blob = { "title": file_name_stem , "slide_objects": [] }
			blob_ulid = str( ULID() )
			blob_file_path = Path( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "local_blob_storage_path" ] ).joinpath( f"{blob_ulid}.json" )
			for index , image_path in enumerate( image_paths ):
				# ulid = str( ULID() )
				# file_path = f"{ulid}.jpeg"
				# token = jwt.encode({ "blob_ulid": blob_ulid } ,
				# 	# utils.base64_decode( config[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ) ,
				# 	config[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ,
				# 	algorithm=config[ "image_upload_server_imgur_version" ][ "jwt_algorithm" ]
				# )
				# print( f"{token.decode( 'utf-8' )} === {ulid}" )
				blob[ "slide_objects" ].append({
					"image_map": image_maps_in_slides[ index ] ,
					# "token": token.decode( "utf-8" ) ,
					"image_ulid": image_ulids[ index ]
				})
			pprint( blob )

			blob_json_str = json.dumps( blob )
			blob_json_b64 = utils.base64_encode( blob_json_str )
			utils.write_json( str( blob_file_path ) , {
				"sealed": utils.secret_box_seal( SECRET_BOX_KEY , blob_json_b64 )
			})
			print( blob_ulid )


			token = jwt.encode({ "blob_ulid": blob_ulid , "key": SECRET_BOX_KEY } ,
				# utils.base64_decode( config[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ) ,
				config[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ,
				algorithm=config[ "image_upload_server_imgur_version" ][ "jwt_algorithm" ]
			)
			token_string = token.decode( "utf-8" )

			print( "\n" )
			drag_and_drop_url = f'{config[ "html"][ "base_hosted_url" ]}/DragAndDrop?t={token_string}'
			typing_url = f'{config[ "html"][ "base_hosted_url" ]}/Typing?t={token_string}'
			print( drag_and_drop_url )
			print( "\n" )
			print( typing_url )

			# print( "here 2" )
			# image_objects = []
			# for slide_index , image_maps_in_slide in enumerate( image_maps_in_slides ):
			# 	print( f"\nSlide === {slide_index+1}" )
			# 	image_objects.append([
			# 		f"Slide" ,
			# 		f"../../images/{image_paths[ slide_index ].name}" ,
			# 		image_maps_in_slide
			# 	])
			# print( "here 3" )
			# pprint( image_objects )

			# # 6.) Generate Interactive Games
			# config[ "html" ][ "base_hosted_url" ] = "./"
			# config[ "html" ][ "cdn" ][ "jquery_ui_css"][ "url" ] = f"../../css/jquery-ui.css"
			# config[ "html" ][ "cdn" ][ "jquery_ui_js"][ "url" ] = f"../../js/jquery-ui.min.js"
			# config[ "html" ][ "cdn" ][ "jquery_js"][ "url" ] = f"../../js/jquery-3.6.0.min.js"
			# config[ "html" ][ "cdn" ][ "bootstrap_css"][ "url" ] = f"../../css/bootstrap.min.css"
			# config[ "html" ][ "cdn" ][ "bootstrap_bundle"][ "url" ] = f"../../js/bootstrap.bundle.min.js"
			# config[ "html" ][ "cdn" ][ "interactive_typing_js" ] = f"../../js/interactive_typing.js"
			# config[ "html" ][ "cdn" ][ "interactive_drag_and_drop_js" ] = f"../../js/interactive_drag_and_drop.js"
			# html_options = config[ "html" ]
			# html_options[ "title" ] = uploaded_powerpoint_path.stem.replace( " " , "-" )
			# html_options[ "images" ] = image_objects
			# html_options[ "output_base_dir" ] = output_html_path
			# interactive_notes_generator.generate( html_options , False )

			# # 7.) Zip Everything and Return
			# shutil.make_archive( temp_dir_posix.joinpath( f"{generated_output_base_path.stem}" ) , "zip" , str( generated_output_base_path ) )
			# output_zip_path = temp_dir_posix.joinpath( f"{generated_output_base_path.stem}.zip" )
			# return await sanic_file(
			# 	str( output_zip_path ) ,
			# 	mime_type="application/zip" ,
			# 	filename=f"{generated_output_base_path.stem}.zip"
			# )
			return sanic_json( dict( testing="unfinished , figuring out where to upload and serve images , with keys and such , routing" ) , status=200 )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )


@app.route( "/test/host/image/<ulid:str>" , methods=[ "GET" ] )
async def local( request: Request , ulid: str ):
	try:
		token = request.args.get( "t" )
		if token == None:
			return sanic_json( dict( failed="no token" ) , status=200 )
		decoded = False
		try:
			decoded = jwt.decode(
				token ,
				# utils.base64_decode( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ) ,
				DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ,
				algorithm=DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "jwt_algorithm" ]
			)
		except Exception as decode_error:
			print( decode_error )
			return sanic_json( dict( failed=str( decode_error ) ) , status=200 )
		if "key" not in decoded:
			return sanic_json( dict( failed="no key sent in token" ) , status=200 )
		encrypted_json_file_path = Path( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "local_image_storage_path" ] ).joinpath( f"{ulid}.json" )
		if encrypted_json_file_path.is_file() == False:
			return sanic_json( dict( failed="file doesn't exist" ) , status=200 )
		encrypted_json = utils.read_json( str( encrypted_json_file_path ) )
		if "sealed" not in encrypted_json:
			return sanic_json( dict( failed="nothing sealed" ) , status=200 )
		opened_base64 = utils.secret_box_open( decoded[ "key" ] , encrypted_json[ "sealed" ] )
		image_bytes = base64.b64decode( opened_base64 )

		# Option 1.) Raw Bytes
		# https://github.com/sanic-org/sanic/blob/main/sanic/response.py#L248
		# return sanic_raw( image_bytes , status=200 , headers={ "Content-Type": "image/jpeg" } )

		# Option 2.) Write Bytes to Temp File and Send
		# with tempfile.TemporaryDirectory() as temp_dir:
		# 	temp_dir_posix = Path( temp_dir )
		# 	with tempfile.NamedTemporaryFile( suffix='.jpeg' , prefix=temp_dir ) as tf:
		# 		temp_file_path = temp_dir_posix.joinpath( tf.name )
		# 		with open( str( temp_file_path ) , 'wb' ) as f:
		# 			f.write( image_bytes )
		# 		return await sanic_file( str( temp_file_path ) )

		# Option 3.) Send Base64 String?
		return sanic_json( dict( image_b64_string=opened_base64 ) , status=200 )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )


@app.route( "/test/host/image" , methods=[ "GET" ] )
async def local( request: Request ):
	try:
		token = request.args.get( "t" )
		if token == None:
			return sanic_json( dict( failed="no token" ) , status=200 )
		decoded = False
		try:
			decoded = jwt.decode(
				token ,
				# utils.base64_decode( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ) ,
				DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ,
				algorithm=DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "jwt_algorithm" ]
			)
		except Exception as decode_error:
			print( decode_error )
			return sanic_json( dict( failed=str( decode_error ) ) , status=200 )
		if "path" not in decoded:
			return sanic_json( dict( failed="no path" ) , status=200 )
		file_path = Path( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "local_image_storage_path" ] ).joinpath( decoded[ "path" ] )
		print( file_path )
		if file_path.is_file() == False:
			return sanic_json( dict( failed="file doesn't exist" ) , status=200 )
		return await sanic_file( str( file_path ) )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )

@app.route( "/test/host/DragAndDrop" , methods=[ "GET" ] )
async def local( request: Request ):
	try:
		token = request.args.get( "t" )
		if token == None:
			return sanic_json( dict( failed="no token" ) , status=200 )
		decoded = False
		try:
			decoded = jwt.decode(
				token ,
				# utils.base64_decode( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ) ,
				DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ,
				algorithm=DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "jwt_algorithm" ]
			)
		except Exception as decode_error:
			print( decode_error )
			return sanic_json( dict( failed=str( decode_error ) ) , status=200 )
		if "blob_ulid" not in decoded:
			return sanic_json( dict( failed="no blob ulid" ) , status=200 )
		if "key" not in decoded:
			return sanic_json( dict( failed="no key" ) , status=200 )
		blob_file_path = Path( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "local_blob_storage_path" ] ).joinpath( f"{decoded[ 'blob_ulid' ]}.json" )
		if blob_file_path.is_file() == False:
			return sanic_json( dict( failed="file doesn't exist" ) , status=200 )
		blob = utils.read_json( str( blob_file_path ) )
		if "sealed" not in blob:
			return sanic_json( dict( failed="nothing sealed ???" ) , status=200 )

		opened_base64 = utils.secret_box_open( decoded[ "key" ] , blob[ "sealed" ] )
		blob = json.loads( base64.b64decode( opened_base64 ) )
		pprint( blob )

		## Now just generate on-the-fly html for drag and drop , and send
		## we are going to have to make a new version of interactive_drag_and_drop.js and interactive_typing.js
		## so that they support the "blob" , and advance and previous work on arrrow keys
		## images get pulled via ulids --> decrypt sealed image base64 string --> render
		html_options = DEFAULT_CONFIG[ "html" ]
		html_options[ "cdn" ][ "jquery_ui_css" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/css/jquery-ui.css"
		html_options[ "cdn" ][ "jquery_ui_js" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/jquery-ui.min.js"
		html_options[ "cdn" ][ "jquery_js" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/jquery-3.6.0.min.js"
		html_options[ "cdn" ][ "bootstrap_css" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/css/bootstrap.min.css"
		html_options[ "cdn" ][ "bootstrap_bundle" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/bootstrap.bundle.min.js"
		html_options[ "cdn" ][ "interactive_typing_js" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/interactive_typing_blob.js"
		html_options[ "cdn" ][ "interactive_drag_and_drop_js" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/interactive_drag_and_drop_blob.js"
		html_options[ "title" ] = blob[ "title" ]
		html_options[ "blob" ] = blob

		html = interactive_notes_generator.build_drag_and_drop_blob_html( html_options )
		return sanic_html( html )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )


@app.route( "/test/host/Typing" , methods=[ "GET" ] )
async def local( request: Request ):
	try:
		token = request.args.get( "t" )
		if token == None:
			return sanic_json( dict( failed="no token" ) , status=200 )
		decoded = False
		try:
			decoded = jwt.decode(
				token ,
				# utils.base64_decode( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ) ,
				DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "sanic_secret_key" ] ,
				algorithm=DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "jwt_algorithm" ]
			)
		except Exception as decode_error:
			print( decode_error )
			return sanic_json( dict( failed=str( decode_error ) ) , status=200 )
		if "blob_ulid" not in decoded:
			return sanic_json( dict( failed="no blob ulid" ) , status=200 )
		if "key" not in decoded:
			return sanic_json( dict( failed="no key" ) , status=200 )
		blob_file_path = Path( DEFAULT_CONFIG[ "image_upload_server_imgur_version" ][ "local_blob_storage_path" ] ).joinpath( f"{decoded[ 'blob_ulid' ]}.json" )
		if blob_file_path.is_file() == False:
			return sanic_json( dict( failed="file doesn't exist" ) , status=200 )
		blob = utils.read_json( str( blob_file_path ) )
		if "sealed" not in blob:
			return sanic_json( dict( failed="nothing sealed ???" ) , status=200 )

		opened_base64 = utils.secret_box_open( decoded[ "key" ] , blob[ "sealed" ] )
		blob = json.loads( base64.b64decode( opened_base64 ) )
		pprint( blob )

		## Now just generate on-the-fly html for drag and drop , and send
		## we are going to have to make a new version of interactive_drag_and_drop.js and interactive_typing.js
		## so that they support the "blob" , and advance and previous work on arrrow keys
		## images get pulled via ulids --> decrypt sealed image base64 string --> render
		html_options = DEFAULT_CONFIG[ "html" ]
		html_options[ "cdn" ][ "jquery_ui_css" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/css/jquery-ui.css"
		html_options[ "cdn" ][ "jquery_ui_js" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/jquery-ui.min.js"
		html_options[ "cdn" ][ "jquery_js" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/jquery-3.6.0.min.js"
		html_options[ "cdn" ][ "bootstrap_css" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/css/bootstrap.min.css"
		html_options[ "cdn" ][ "bootstrap_bundle" ][ "url" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/bootstrap.bundle.min.js"
		html_options[ "cdn" ][ "interactive_typing_js" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/interactive_typing_blob.js"
		html_options[ "cdn" ][ "interactive_drag_and_drop_js" ] = f"{html_options[ 'base_hosted_url' ]}/static/js/interactive_drag_and_drop_blob.js"
		html_options[ "title" ] = blob[ "title" ]
		html_options[ "blob" ] = blob

		html = interactive_notes_generator.build_typing_blob_html( html_options )
		return sanic_html( html )
	except Exception as e:
		print( e )
		return sanic_json( dict( failed=str( e ) ) , status=200 )

if __name__ == "__main__":
	app.run( host="0.0.0.0" , port="9379" )