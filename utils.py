import sys
import json
import base64
from pathlib import Path
from natsort import humansorted

import nacl.secret
import nacl.utils

from pptx import Presentation
from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches , Pt
from pptx.oxml import parse_xml
from pptx.dml.color import _NoneColor , RGBColor

def base64_encode( message ):
	try:
		message_bytes = message.encode( 'utf-8' )
		base64_bytes = base64.b64encode( message_bytes )
		base64_message = base64_bytes.decode( 'utf-8' )
		return base64_message
	except Exception as e:
		print( e )
		return False

def base64_decode( base64_message ):
	try:
		base64_bytes = base64_message.encode( 'utf-8' )
		message_bytes = base64.b64decode(base64_bytes)
		message = message_bytes.decode( 'utf-8' )
		return message
	except Exception as e:
		print( e )
		return False

def write_json( file_path , python_object ):
	with open( file_path , 'w', encoding='utf-8' ) as f:
		json.dump( python_object , f , ensure_ascii=False , indent=4 )

def read_json( file_path ):
	with open( file_path ) as f:
		return json.load( f )

# https://stackoverflow.com/a/24290026
def enumerate2( xs , start=0 , step=1 ):
	for x in xs:
		yield ( start , x )
		start += step

def secret_box_generate_new_key():
	key = nacl.utils.random( nacl.secret.SecretBox.KEY_SIZE )
	key_b64 = base64.b64encode( key ).decode( "utf-8" )
	print( f"Human Readable Key \t=== {key_b64}" )

def secret_box_seal( key_base64 , plain_text_message ):
	box = nacl.secret.SecretBox( base64.b64decode( key_base64 ) )
	encrypted = box.encrypt( bytes( plain_text_message , "utf-8" ) )
	encrypted_b64 = base64.b64encode( encrypted ).decode( "utf-8" )
	return encrypted_b64

def secret_box_open( key_b64 , encrypted_message ):
	key_bytes = base64.b64decode( key_b64 )
	encrypted_message_bytes = base64.b64decode( encrypted_message )
	box = nacl.secret.SecretBox( key_bytes )
	plaintext = box.decrypt( encrypted_message_bytes ).decode( "utf-8" )
	return plaintext

# https://stackoverflow.com/a/49111747
def get_nested_dictionary_value( d , l , default_val=None ):
	if l[ 0 ] not in d:
		print( "1" )
		return default_val
	elif len( l )==1:
		print( "2" )
		return d[ l[ 0 ] ]
	else:
		print( "3" )
		return get_nested_dictionary_value( d[ l[ 0 ] ] , l[ 1: ] )

def set_nested_dictionary_value( d , l , set_value=None ):
	if l[ 0 ] not in d:
		print( "1" )
		return set_value
	elif len( l )==1:
		print( "3" )
		d[ l[ 0 ] ] = set_value
		return d
	else:
		print( "3" )
		return set_nested_dictionary_value( d[ l[ 0 ] ] , l[ 1: ] , set_value )

def get_slide_image_paths( power_point_file ):
	power_point_file_posix = Path( power_point_file )
	ALLOWED_EXTENSIONS = [ ".jpeg" ]
	BaseDirectoryPosixPath = power_point_file_posix.parent.joinpath( power_point_file_posix.stem )
	FilesPosixInBaseDirectory = BaseDirectoryPosixPath.glob( '*' )
	FilesPosixInBaseDirectory = [ x for x in FilesPosixInBaseDirectory if x.is_file() ]
	FilesPosixInBaseDirectory = [ x for x in FilesPosixInBaseDirectory if x.suffix in ALLOWED_EXTENSIONS ]
	# FilesPosixInBaseDirectory = humansorted( FilesPosixInBaseDirectory )
	FilesPosixInBaseDirectory = [ BaseDirectoryPosixPath.joinpath( f"Slide{x}.jpeg" ) for x in range( 1 , len( FilesPosixInBaseDirectory ) + 1 ) ]
	return FilesPosixInBaseDirectory


def is_rectangular( shape ):
	if shape.width > shape.height:
		return True
	return False

def is_accepted_rectangular_form( shape ):
	global MAXIMUM_TEXTBOX_HEIGHT
	# if is_rectangular( shape ) == False:
	# 	return False
	# if shape.height > MAXIMUM_TEXTBOX_HEIGHT:
	# 	return False
	return True

def get_fill_type( fill ):
	if hasattr( fill.type , "real" ) == False:
		return False
	fill_type = fill.type
	if fill_type == 1:
		return "solid"
	elif fill_type == 2:
		return "patterned"
	elif fill_type == 3:
		return "gradient"
	elif fill_type == 4:
		return "textured"
	elif fill_type == 5:
		return "background"
	elif fill_type == 6:
		return "picture"
	return False

def get_fill_color( shape ):
	if hasattr( shape , "fill" ) == False:
		return False
	fill_type = get_fill_type( shape.fill )
	if fill_type == False:
		return False
	if fill_type != "solid":
		return False
	color = shape.fill.fore_color
	result = {}
	if not isinstance( color._color , _NoneColor ):
		if hasattr( color, "rgb" ):
			result[ "rgb" ] = color.rgb
		if hasattr( color , "brightness" ):
			result[ "brightness" ] = color.brightness
		if hasattr( color , "color_type" ):
			result[ "color_type" ] = color.color_type
		if hasattr( color , "theme_color" ):
			result[ "theme_color" ] = color.theme_color
		if hasattr( color , "theme_color" ):
			result[ "theme_color" ] = "{}".format( color.theme_color )
	return result

def validate_correct_color( rgb_object , valid_hex_color ):
	if str( rgb_object ) == valid_hex_color:
		return True
	# elif isinstance( rgb_object , tuple ):
	# 	if rgb_object[ 0 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 0 ]:
	# 		if rgb_object[ 1 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 1 ]:
	# 			if rgb_object[ 2 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 2 ]:
	# 				return True
	return False

def validate_text_box( shape  , valid_hex_color ):
	# 1.) Validate Background Color Matches "our" Textboxes
	# NOTE: YOU HAVE TO USE SHAPE FILL BUTTON , idk man
	fill_color = get_fill_color( shape )
	if fill_color == False:
		return False
	if "rgb" not in fill_color:
		return False
	if validate_correct_color( fill_color[ "rgb" ] , valid_hex_color ) == False:
		return False

	# 2.) Validate There is Text In the Box
	if hasattr( shape , "text_frame" ) == False:
		return False
	if len( shape.text_frame.text ) < 1:
		return False
	if shape.text_frame.text == " ":
		return False

	# 3.) Validate Shape is "rEcTanGULar" and appropriate size
	# doesn't do a very good job, but 99% it works
	if is_accepted_rectangular_form( shape ) == False:
		return False
	return True