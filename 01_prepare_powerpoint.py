#!/usr/bin/env python3
import sys
from pathlib import Path
from pprint import pprint
from pptx import Presentation
from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches , Pt
from pptx.oxml import parse_xml
from pptx.dml.color import _NoneColor , RGBColor

import utils
import compute_image_maps

OUR_BACKGROUND_TEXTBOX_HEX = "0070C0"
OUR_BACKGROUND_TEXTBOX_RGB = [ 68 , 114 , 196 ]

def is_rectangular( shape ):
	if shape.width > shape.height:
		return True
	return False

def is_accepted_rectangular_form( shape ):
	global MAXIMUM_TEXTBOX_HEIGHT
	# if is_rectangular( shape ) == False:
	# 	return False
	# if shape.height > MAXIMUM_TEXTBOX_HEIGHT:
	# 	return False
	return True

def get_fill_type( fill ):
	if hasattr( fill.type , "real" ) == False:
		return False
	fill_type = fill.type
	if fill_type == 1:
		return "solid"
	elif fill_type == 2:
		return "patterned"
	elif fill_type == 3:
		return "gradient"
	elif fill_type == 4:
		return "textured"
	elif fill_type == 5:
		return "background"
	elif fill_type == 6:
		return "picture"
	return False

def get_fill_color( shape ):
	if hasattr( shape , "fill" ) == False:
		return False
	fill_type = get_fill_type( shape.fill )
	if fill_type == False:
		return False
	if fill_type != "solid":
		return False
	color = shape.fill.fore_color
	result = {}
	if not isinstance( color._color , _NoneColor ):
		if hasattr( color, "rgb" ):
			result[ "rgb" ] = color.rgb
		if hasattr( color , "brightness" ):
			result[ "brightness" ] = color.brightness
		if hasattr( color , "color_type" ):
			result[ "color_type" ] = color.color_type
		if hasattr( color , "theme_color" ):
			result[ "theme_color" ] = color.theme_color
		if hasattr( color , "theme_color" ):
			result[ "theme_color" ] = "{}".format( color.theme_color )
	return result

def validate_correct_color( rgb_object ):
	global OUR_BACKGROUND_TEXTBOX_RGB
	global OUR_BACKGROUND_TEXTBOX_HEX
	if str( rgb_object ) == OUR_BACKGROUND_TEXTBOX_HEX:
		return True
	elif isinstance( rgb_object , tuple ):
		if rgb_object[ 0 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 0 ]:
			if rgb_object[ 1 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 1 ]:
				if rgb_object[ 2 ] == OUR_BACKGROUND_TEXTBOX_RGB[ 2 ]:
					return True
	return False

def validate_text_box( shape ):
	# 1.) Validate Background Color Matches "our" Textboxes
	# NOTE: YOU HAVE TO USE SHAPE FILL BUTTON , idk man
	fill_color = get_fill_color( shape )
	if fill_color == False:
		return False
	if "rgb" not in fill_color:
		return False
	if validate_correct_color( fill_color[ "rgb" ] ) == False:
		return False

	# 2.) Validate There is Text In the Box
	if hasattr( shape , "text_frame" ) == False:
		return False
	if len( shape.text_frame.text ) < 1:
		return False
	if shape.text_frame.text == " ":
		return False

	# 3.) Validate Shape is "rEcTanGULar" and appropriate size
	# doesn't do a very good job, but 99% it works
	if is_accepted_rectangular_form( shape ) == False:
		return False
	return True

# def generate_local_stage_one( input_path , config , output_path ):
# 	input_path = Path( sys.argv[ 1 ] )
# 	config = utils.read_json( sys.argv[ 2 ] )

# 	# 1.) Create Copy of PowerPoint With All text removed from boxes with correct fill color
# 	p = Presentation( sys.argv[ 1 ] )
# 	p_clone = deepcopy( p )
# 	for slide_index , slide in enumerate( p_clone.slides ):
# 		for shape_index , shape in enumerate( slide.shapes ):
# 			if validate_text_box( shape ):
# 				print( f"{slide_index} === {shape_index} === valid text box" )
# 				shape.text_frame.text = ""
# 	p_clone.save( f"{input_path.stem}-Blank.pptx" )


if __name__ == "__main__":

	input_path = Path( sys.argv[ 1 ] )
	config = utils.read_json( sys.argv[ 2 ] )
	if "parser" in config:
		if "our_background_textbox_hex" in config[ "parser" ]:
			OUR_BACKGROUND_TEXTBOX_HEX = config[ "parser" ][ "our_background_textbox_hex" ]
		if "our_background_textbox_rgb" in config[ "parser" ]:
			OUR_BACKGROUND_TEXTBOX_RGB = config[ "parser" ][ "our_background_textbox_rgb" ]

	# 1.) Create Copy of PowerPoint With All text removed from boxes with correct fill color
	p = Presentation( sys.argv[ 1 ] )
	p_clone = deepcopy( p )
	for slide_index , slide in enumerate( p_clone.slides ):
		for shape_index , shape in enumerate( slide.shapes ):
			if validate_text_box( shape ):
				print( f"{slide_index} === {shape_index} === valid text box" )
				shape.text_frame.text = ""
	p_clone.save( f"{input_path.stem}-Blank.pptx" )

	# # 2.) Just Scan Image Maps Now so we never need original file again
	# image_maps_in_slides = compute_image_maps.compute( sys.argv[ 1 ] , config[ "parser" ] )
	# # pprint( image_maps_in_slides )
	# utils.write_json( f"{input_path.stem}-image-maps.json" , image_maps_in_slides )